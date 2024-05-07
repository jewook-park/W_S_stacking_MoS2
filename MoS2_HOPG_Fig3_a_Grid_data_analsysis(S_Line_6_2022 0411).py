# -*- coding: utf-8 -*-
# ---
# jupyter:
#   jupytext:
#     formats: ipynb,py:light
#     text_representation:
#       extension: .py
#       format_name: light
#       format_version: '1.5'
#       jupytext_version: 1.13.7
#   kernelspec:
#     display_name: Python 3 (ipykernel)
#     language: python
#     name: python3
# ---

# + [markdown] id="P78u1O1Kblyo" jp-MarkdownHeadingCollapsed=true tags=[] jp-MarkdownHeadingCollapsed=true
# #  $MoS_{2}$ islands on HOPG substrates
# # STM  data analysis data analysis (Fig 3- layer-dependent STS  of S islands)
# # based on **SPMpy** data file loading platform 
#
# # SPMPY data loading & analysis (3D data, *.3ds) 
#
# >> SPMPY aims for a python package for code-based scanning probe microscopy data analysis. Scanning probe microscopy(SPM), such as scanning tunneling microscopy and spectroscopy (STM/S) and atomic force microscopy (AFM), generate an inherently multidimensional dataset, which possesses scientific information in nanoscale. SPMPY exploits recent image processing techniques, a.k.a. computer vision technology, to the multidimensional dataset. Automated data processing with computational power will extract scientific information from experimental data and accelerate cutting-edge research in nanoscience.
#
# * SPMPY provides data loading platform for the Nanonis data set (* .dat, *.sxm, *.3ds) to convert the multidimensional SPM data as PANDAS DataFrames. To take advantage of SPMPY, other types of SPM data are supposed to be converted as DataFrame in advance. 
#
# > Author: Dr. Jewook Park  (jewookpark@ibs.re.kr)
#  IBS-CALDES (Inistitute for Basic Science, Center for Artificial Low Dimensional Electronic Systems), South Korea
#
# * Version 0.2022.0208
#
#
# ## Goal
# * Analyze 3D SPM image 
#     *Especially,  **Nanonis** Grid data (*.3ds) file loading & visualization
# * sxm data files ↦  **Xarray** (DataSet or DataArray) (or **PANDAS**  (DataFrame) )
#
# > This notebook is intended as an example of **SPMpy** project.
#  **SPMpy** is an open source project initiated by Dr.Jewook Park@IBS-CALDES. 
# Github : https://github.com/Jewook-Park/SPMPY
# * SPMpy is on-going project. 
# * Contribution, comments, ideas, error reports are always welcome.(Korean or English)
#
#
#
#
# ## pptx file save => new function  (to be added)
# * save all images in a pptx file with scan info (description)
#

# + [markdown] tags=[]
# # Import modules and functions 
#
# ## **1.Preparation** 
# * Import necessary packages
#     
# ### loading SPMpy functions
#     
# * SPMpy_file_loading_funcs.py 
# * SPMpy_data_analysis_funcs_sxm.py 
#

# + id="Qm1zLaTHbSpK"
########################################
    #    * Step 1-1
    #    : Import necessary packages 
    #        import modules        
#########################################

import os
import glob
import numpy as np
import pandas as pd
import scipy as sp
from warnings import warn
from scipy import signal

import math
import matplotlib.pyplot as plt
import matplotlib.patches as patches

import seaborn as sns
import skimage
import seaborn_image as isns

import panel as pn
import panel.widgets as pnw
import ipywidgets as ipw

from scipy.signal import find_peaks

# some packages may be yet to be installed
# please install "seaborn-image" via conda to avoid install error 
# conda install -c conda-forge seaborn-image
try:
     from pptx import Presentation
except ModuleNotFoundError:
    warn('ModuleNotFoundError: No module named Presentation')
    # !pip install python-pptx  
    from pptx import Presentation
    from pptx.util import Inches, Pt

try:
    import nanonispy as nap
except ModuleNotFoundError:
    warn('ModuleNotFoundError: No module named nanonispy')
    # !pip install nanonispy
    import nanonispy as nap



try:
    import xarray as xr
except ModuleNotFoundError:
    warn('ModuleNotFoundError: No module named xarray')
    
    # !pip install xarray 
    import xarray as xr
    
try:
    import xrft
except ModuleNotFoundError:
    warn('ModuleNotFoundError: No module named xrft')
    # !pip install xrft 
    import xrft
    
    
try:
    import holoviews as hv
except ModuleNotFoundError:
    warn('ModuleNotFoundError: No module named holoviews')
    # !pip install holoviews 
    import holoviews as hv
    
try:
    import hvplot.xarray
    import hvplot.pandas 
except ModuleNotFoundError:
    warn('ModuleNotFoundError: No module named hvplot')
    # !pip install hvplot
    import hvplot.xarray
    import hvplot.pandas 

    
    
from SPMpy_file_loading_funcs import *
from SPMpy_2D_data_analysis_funcs import *
from SPMpy_3D_data_analysis_funcs import *




# + [markdown] id="gzPbBV16fOU6" tags=[]
# # 2. Grid DATA Import 
#     
# * check all SPM data files in the given folder 
#     * use **FilesInFolder** function
#         * input : **file path**
#         * return DataFrame columns = ['group','num','filename']
# * use the **Grid2xr** function to load the sxm file as a xarray format
#

# +
# check the sxm files in the given folder
#target_path = r'D:\CALDES data - JEWOOK PARK research group\CloudStation\SPMs\ULT-SPM (Unisoku) - Jewook Park\Raw data\2021\2021 1229 Wtip16 Gr_SiC LN2T'

target_path = r'D:\CALDES data - JEWOOK PARK research group\CloudStation\SPMs\LT STM (Unisoku) - Yeom\Raw data (2016-p)\2018 0314- MoS2_HOPG_Yonsei 78.15K'
#target_path = r'C:\CloudStation\CloudStation-  ULT-SPM\CloudStation - HOPH _ 0711'

## ###
# add 'r' to avoid unicodeerror 
#######
file_list_df = FilesInFolder(target_path)

grid_file_list_df = file_list_df[file_list_df.type=='3ds']
# file_list_df
grid_file_groups = list (set(grid_file_list_df['group']))
# file groups
grid_file_groups

#############################
# 3ds file loading by using Grid2xr
# xr format으로 파일 불러오기 
# (e.g.) # Grid2xr(file_list_df.file_name.iloc[0])
#############################
#spmdata_xr = img2xr(file_list_df.file_name[0])



# + [markdown] tags=[]
# # Choose the grid data file to analyze
#
# * using **Grid2xr** function 
#     * add fwd/bwd data average '_fb'
#     * split 2D topo info & 3D info in the grid 
#         * grid_topo : 2D data 
#         * grid_3D : 3D data
#     * **I_fb** : I, (forwad + backward sweep )/2
#     * **LIX_fb** : LIX, (forwad + backward sweep )/2
#     * **dIdV** : dI/dV (using xr differentiate _class_ )
#     * **LIX_unit_calc** : LIX_fb *  LIX_coefficient (for unit calibration) 
#
# * after **grid_3D_Gap** function * after **grid_3D_Gap** function 
#     *  2D data :  CBM, VBM position assignment $\leftarrow$ based on **I** or **LIX**
#         * CBM_I_mV, VBM_I_mV, gap_size_I 
#         * CBM_LIX_mV, VBM_LIX_mV, gap_size_LIX
#     *  3D data : LDOS_fb $\leftarrow$ after unit calc & offset adjust
#         * I_fb, LIX_fb, LDOS_fb, LDOS_fb_CB, LDOS_fb_VB
#         
#         * **I_fb** : I, (forwad + backward sweep )/2
#         * **LIX_fb** : LIX, (forwad + backward sweep )/2
#         * **LDOS_fb** : LIX_fb *  LIX_coefficient (for unit calibration)  + offset adjustment (according to LIX  at I=0)
#         * **LDOS_fb_CB** : based on LIX assignment
#         * **LDOS_fb_VB** : based on LIX assignment
#     
#     
# -

#grid_xr = Grid2xr(grid_file_list_df.file_name.iloc[0])
# SELECT  data 
Spiral_line_006 = grid_file_list_df[grid_file_list_df.file_name == "Grid Spectroscopy_line006.3ds"]
Spiral_line_006.file_name.values[0]

grid_xr = Grid2xr(Spiral_line_006.file_name.values[0])

grid_xr

# # imperfect grid data 
#
# ## select non-zero are 

grid_xr.where(grid_xr.topography != 0).topography.plot()

# +
slct_Y = grid_xr.where(grid_xr.topography != 0).dropna(dim = 'Y').Y
# use the dropna  to select Y range 
grid_xr = grid_xr.sel(Y = slct_Y[:-4])



#grid_xr
# -

slct_Y = grid_xr.where(grid_xr.topography != 0).dropna(dim = 'Y').Y
slct_Y

grid_xr

# + jp-MarkdownHeadingCollapsed=true tags=[]
# grid data to xr 
grid_xr['I_fb'] = (grid_xr.I_fwd + grid_xr.I_fwd)/2
grid_xr['LIX_fb'] = (grid_xr.LIX_fwd + grid_xr.LIX_fwd)/2
# add 'I' & 'LIX' channel "fb = [fwb+bwd] / 2 " 

grid_topo = grid_xr[['topography']]
# topography 
grid_3D = grid_xr[['I_fb','LIX_fb']]
# averaged I & LIX 
#grid_topo.topography.plot()
#grid_3D.LIX_fb.isel(bias_mV=0).plot()

grid_3D_gap = grid_3D_Gap(grid_3D, I_0_pA=1E-14, LIX_0_pA=1E-15)
# assign gap from STS
#grid_3D_gap

# +
grid_topo.where(grid_topo.topography != 0).topography.plot()

#grid_3D

## LDOS unit calibration 

# + [markdown] tags=[]
# # Pretreatments 
#
# ## Find Energy gap in STS 
# * (I or LIX) is almost Zero. less than measurment error bar
#     * I_min_pA= 1E-11 // LIX_min_pA= 1E-12
#     * find gap size & adjust LIX offset, based on I or LIX
#         * to prevent ripples in dI/dV polyfit 
#

# + [markdown] tags=[] jp-MarkdownHeadingCollapsed=true
# # interfactive plot with hvplot
#
# -

grid_topo

from hvplot import hvPlot
import holoviews as hv
import hvplot.xarray  # noqa
import panel as pn
import panel.widgets as pnw
import ipywidgets as ipw
from holoviews import opts
from holoviews.streams import Stream, param


# + [markdown] tags=[]
# # 1. 3D plot & slicing 
#
# * use the holoview 
#     * or hvplot(with panel- widget) interactive plot  $\to$  event handling  $\to$  later 
#     * mayavi? $\to$ later
#

# +
#grid_topo.topography.plot()
#isns.imshow(grid_topo.topography)
# Plane fit (avg = 0 )
#PlaneFitY_xr(grid_topo).min().topography.plot()

#set min = 0 
z_profile  = (PlaneFitY_xr(grid_topo) - PlaneFitY_xr(grid_topo).min()).to_dataframe().reset_index()
#z_profile
z_profile = z_profile* 1E9
sns.set_style("whitegrid")
z_profile_plot = sns.lineplot(x = "X", y = "topography", data = z_profile)
z_profile_plot.set_xlabel ("X (nm)")
z_profile_plot.set_ylabel ("Z (nm)")
#z_profile.topography.plot()

# -

# ## xr to hv 

# +
###############
# holoview slicing 
###############

# choose the channel []

grid_3D_hv = hv.Dataset(grid_3D_gap.LIX_fb)
# convert xr dataset as a holoview dataset 
hv.extension('bokeh')
# -

# ##  Bias_mV slicing Examples 

###############
# bias_mV slicing
dmap_plane  = ["X","Y"]
dmap = grid_3D_hv.to(hv.Image,
                     kdims = dmap_plane,
                     dynamic = True )
dmap.opts(colorbar = True,
          cmap = 'bwr',
          frame_width = 400,
          aspect = 1).relabel('XY plane slicing: ')
fig = hv.render(dmap)
dmap   

# + [markdown] jp-MarkdownHeadingCollapsed=true tags=[] jp-MarkdownHeadingCollapsed=true jp-MarkdownHeadingCollapsed=true tags=[]
# ##  Y & X slicing Examples 
# -

###############
# Y slicing
dmap_plane  = [ "X","bias_mV"]
dmap = grid_3D_hv.to(hv.Image,
                     kdims = dmap_plane,
                     dynamic = True )
dmap.opts(colorbar = True,
          cmap = 'bwr',
          frame_width = 400).relabel('X - bias_mV plane slicing: ')

# +
###############
# X slicing
dmap_plane  = ["bias_mV","Y"]
dmap = grid_3D_hv.to(hv.Image,
                     kdims = dmap_plane,
                     dynamic = True )
dmap.opts(colorbar = True,
          cmap = 'bwr',
          frame_width = 400).relabel('Y - bias_mV plane slicing: ')

# key dimension assign 
# dynamic = True ==> use the DynamicMap, False ==> use the HoloMap (which uses more memory)
# or use the  'fig_size=200'(matplotlib?) or 'frame_width = 400'(bokeh?) for image sizes 

# setting color ranges, normalization, etc.. .
# http://holoviews.org/user_guide/Style_Mapping.html
# -

# ## Terrace check 

# +
z_profile  = (PlaneFitY_xr(grid_topo) - PlaneFitY_xr(grid_topo).min())
z_profile

# check the terrrace number by using MultiOtsu
z_profile_layer = ThrshldMultiOtsu_xr(z_profile)
z_profile_layer = z_profile_layer.topography.rename('MOtsu')

z_profile_layer

#plot test after Multi Ots
z_profile = z_profile.merge(z_profile_layer)



# +
z_profile_layer.plot()

z_profile_layer
# -

z_profile_sobel = FilterSobel_xr(z_profile)
z_profile_sobel.MOtsu_sobel.plot()

# +
z_profile_sobel_df = z_profile_sobel.to_dataframe()


mask_terrace = ( z_profile_sobel_df.MOtsu_sobel == 0 )
mask_2ML = (z_profile_sobel_df.MOtsu == 0) & mask_terrace
mask_3ML = (z_profile_sobel_df.MOtsu == 1) & mask_terrace
mask_4ML = (z_profile_sobel_df.MOtsu == 2) & mask_terrace

z_profile_sobel_df['layer'] = 'boundary'
z_profile_sobel_df.layer[mask_2ML]='2ML'
z_profile_sobel_df.layer[mask_3ML]='3ML'
z_profile_sobel_df.layer[mask_4ML]='4ML'

layer_N = z_profile_sobel_df.to_xarray()
# MultiOtsu 활용한 결과를 xarray로 바꿈. 
grid_3D['layer'] = layer_N.layer
grid_3D
# -

z_profile_sobel_df.MOtsu.plot()

grid_3D_gap = grid_3D_Gap(grid_3D, I_0_pA=1E-11, LIX_0_pA= 1E-15)
grid_3D_gap


grid_3D_gap.LDOS_fb.isel(bias_mV=0).plot()

grid_3D_gap_sg = savgolFilter_xr(grid_3D_gap, polyorder= 3, window_length= 11)
grid_3D_gap_sg

grid_3D_gap_sg_1deriv =  grid_3D_gap_sg.differentiate(coord = 'bias_mV')
grid_3D_gap_sg_2deriv = savgolFilter_xr(grid_3D_gap_sg_1deriv).differentiate(coord = 'bias_mV')
grid_3D_gap_sg_2deriv_sg_dps = find_peaks_xr(-1*savgolFilter_xr(grid_3D_gap_sg_2deriv))
grid_3D_gap_sg_2deriv_sg_dps_pad = peak_pad(grid_3D_gap_sg_2deriv_sg_dps)
grid_3D_gap_sg_2deriv_sg_dps_pad

# +
#lineSTS_P = grid_3D_gap_sg[['LDOS_fb','layer']].to_dataframe()
#lineSTS_P = grid_3D_gap_sg[['LDOS_fb','layer']].to_dataframe().reset_index()
# after gap cal, LDOS fb for dI/dV plot 


grid_3D_gap_STS = xr.merge(
    [grid_3D_gap_sg[['I_fb','LDOS_fb','CBM_I_mV','VBM_I_mV','CBM_LIX_mV','VBM_LIX_mV','layer']] ,
     grid_3D_gap_sg_2deriv_sg_dps_pad[['LDOS_fb_peaks']]]).dropna(dim = 'X')

layer_STS_P = grid_3D_gap_STS.groupby('layer').mean(dim='stacked_X_Y')
layer_STS_P
# how select layers in xarray
#layer_STS_P[dict (layer = [0])]

STS_P_2ML = layer_STS_P.sel(layer = '2ML').LDOS_fb
STS_P_3ML = layer_STS_P.sel(layer = '3ML').LDOS_fb
STS_P_4ML = layer_STS_P.sel(layer = '4ML').LDOS_fb

STS_P_2ML.plot()
STS_P_3ML.plot()
STS_P_4ML.plot()
# -

# # use the 2 3 4 ML assign  as calculated. 
#
# ## use the averaged curve to extract the peak position 
#
#
#

STS_P_2ML.to_dataframe().drop('layer', axis =1).rename({'LDOS_fb':'S_2ML'})

# +
S_2ML_df = STS_P_2ML.to_dataframe().drop('layer', axis =1)
S_3ML_df = STS_P_3ML.to_dataframe().drop('layer', axis =1)
S_4ML_df = STS_P_4ML.to_dataframe().drop('layer', axis =1)

S_234ML_df = pd.concat ([S_2ML_df,S_3ML_df,S_4ML_df], axis =1)
S_234ML_df.columns =  ['S_2ML','S_3ML','S_4ML']
# column name chnage 
# -



#sns.lineplot(data = S_2ML_df, x = 'bias_mV', y = 'LDOS_fb')
sns.lineplot(data = S_234ML_df)

# +
#(grid_3D_gap_sg.layer !='boundary').values

grid_3D_gap_sg_drop = grid_3D_gap_sg.where(grid_3D_gap_sg.layer !='boundary')
####################
# remove boundary region 

# +
S_STS_234 = grid_3D_gap_sg_drop[['LDOS_fb','layer']].to_dataframe()

fig, axs = plt.subplots()
sns.lineplot(data = S_STS_234, x = 'bias_mV', y = 'LDOS_fb', hue = 'layer', ax = axs)
axs.set_xlim(-2100, -1200)
#axs.set_ylim(None, 1E-10)
axs.set_yscale('log')
plt.show()
# -

# * find  curve data from S_STS_234

# # Use the find peaks of S_STS_234
#
#
#

#S_234ML_df
S_234_sns = sns.lineplot(data = S_STS_234, x = 'bias_mV', y = 'LDOS_fb', hue = 'layer')
#S_234_sns.get_lines()[0]

# +
#S_234_sns.get_lines()[0].get_data()
S_2ML_df = pd.DataFrame(S_234_sns.get_lines()[2].get_data()).T
S_3ML_df = pd.DataFrame(S_234_sns.get_lines()[1].get_data()).T
S_4ML_df = pd.DataFrame(S_234_sns.get_lines()[0].get_data()).T
# 2,3,4 순서에 유의 
# 위에서 4,3,2 순서로 불러서 그렸기 때문에 순서 바꿔야함. 


S_2ML_df.columns = ['bias_mV','S_2ML']
S_3ML_df.columns = ['bias_mV','S_3ML']
S_4ML_df.columns = ['bias_mV','S_4ML']

sns.lineplot(data = S_2ML_df, x ='bias_mV', y = 'S_2ML')
sns.lineplot(data = S_3ML_df, x ='bias_mV', y = 'S_3ML')
sns.lineplot(data = S_4ML_df, x ='bias_mV', y = 'S_4ML')

# +

#S_234ML_df= pd.concat([S_2ML,S_3ML.S_3ML_LDOS,S_4ML.S_4ML_LDOS], axis =1).set_index('bias_mV')
#S_234ML_df= S_234ML_df.set_index('bias_mV')
S_234ML_xr = pd.concat([S_2ML_df,S_3ML_df.S_3ML,S_4ML_df.S_4ML], axis =1).set_index('bias_mV').to_xarray()
#S_234ML_xr


# +
from scipy.signal import savgol_filter
# 1st derivative 
'''
W_234ML_xr_1deriv = W_234ML_xr.differentiate(coord = 'bias_mV')
W_234ML_xr_1deriv_sg = W_234ML_xr_1deriv.copy()
'''
S_234ML_xr_1deriv = S_234ML_xr.differentiate(coord = 'bias_mV')
S_234ML_xr_1deriv_sg = S_234ML_xr_1deriv.copy()

'''
for ch in W_234ML_xr_1deriv:
    W_234ML_xr_1deriv_sg[ch].values =  savgol_filter(W_234ML_xr_1deriv[ch].values,
                                                     window_length = 7,
                                                     polyorder = 3)
'''   
for ch in S_234ML_xr_1deriv:
    S_234ML_xr_1deriv_sg[ch].values =  savgol_filter(S_234ML_xr_1deriv[ch].values,
                                                     window_length = 5,
                                                     polyorder = 3)
# -

S_234ML_xr_1deriv

#sns.lineplot(data = W_234ML_xr_1deriv_sg.to_dataframe())
sns.lineplot(data = S_234ML_xr_1deriv_sg.to_dataframe())

# +
# 2nd derivative 
'''
W_234ML_xr_2deriv = W_234ML_xr_1deriv_sg.differentiate(coord = 'bias_mV')
W_234ML_xr_2deriv_sg = W_234ML_xr_2deriv.copy()
'''
S_234ML_xr_2deriv = S_234ML_xr_1deriv_sg.differentiate(coord = 'bias_mV')
S_234ML_xr_2deriv_sg = S_234ML_xr_2deriv.copy()
'''
for ch in W_234ML_xr_2deriv:
    W_234ML_xr_2deriv_sg[ch].values =  savgol_filter(W_234ML_xr_2deriv[ch].values,
                                                     window_length = 7,
                                                     polyorder = 3)
'''
for ch in S_234ML_xr_2deriv:
    S_234ML_xr_2deriv_sg[ch].values =  savgol_filter(S_234ML_xr_2deriv[ch].values,
                                                     window_length = 5,
                                                     polyorder = 3)
    


# -

#sns.lineplot(data = W_234ML_xr_2deriv_sg.to_dataframe())
sns.lineplot(data = S_234ML_xr_2deriv_sg.to_dataframe())

# find the 1st derivative peak  to assgin band edge 
'''
for ch in W_234ML_xr_1deriv_sg:
    W_234ML_xr_1deriv_sg[ch+'_pks'] = xr.DataArray(find_peaks(W_234ML_xr_1deriv_sg[ch].values))
    W_234ML_xr_1deriv_sg[ch+'_dps'] = xr.DataArray(find_peaks(-1*W_234ML_xr_1deriv_sg[ch].values))
    # channel 별로 서로다른 number  of peaks  --> 하나로 묶어야함. 
    # 하나로 묶어내고 나서 값을 불러오는 방법은? 
    ############################################
    # W_234ML_xr_2deriv_sg.W_3ML_dps.data.item()[0]
    ############################################
'''    
for ch in S_234ML_xr_1deriv_sg:
    S_234ML_xr_1deriv_sg[ch+'_pks'] = xr.DataArray(find_peaks(S_234ML_xr_1deriv_sg[ch].values))
    S_234ML_xr_1deriv_sg[ch+'_dps'] = xr.DataArray(find_peaks(-1*S_234ML_xr_1deriv_sg[ch].values))



# find the 2st derivative peak   to assgin peak position in dIdV
'''
for ch in W_234ML_xr_2deriv_sg:
    W_234ML_xr_2deriv_sg[ch+'_pks'] = xr.DataArray(find_peaks(W_234ML_xr_2deriv_sg[ch].values))    
    W_234ML_xr_2deriv_sg[ch+'_dps'] = xr.DataArray(find_peaks(-1*W_234ML_xr_2deriv_sg[ch].values))
    # channel 별로 서로다른 number  of peaks  --> 하나로 묶어야함. 
    # 하나로 묶어내고 나서 값을 불러오는 방법은? 
    ############################################
    # W_234ML_xr_2deriv_sg.W_3ML_dps.data.item()[0]
    ############################################
''' 
for ch in S_234ML_xr_2deriv_sg:
    S_234ML_xr_2deriv_sg[ch+'_pks'] = xr.DataArray(find_peaks(S_234ML_xr_2deriv_sg[ch].values))
    S_234ML_xr_2deriv_sg[ch+'_dps'] = xr.DataArray(find_peaks(-1*S_234ML_xr_2deriv_sg[ch].values))

S_234ML_xr_2deriv_sg

# +
#W_234ML_xr_1deriv_sg
# pks N dps
'''
W1_2_pks = W_234ML_xr_1deriv_sg.W_2ML_pks.data.item()[0]
W1_3_pks = W_234ML_xr_1deriv_sg.W_3ML_pks.data.item()[0]
W1_4_pks = W_234ML_xr_1deriv_sg.W_4ML_pks.data.item()[0]
W1_2_dps = W_234ML_xr_1deriv_sg.W_2ML_dps.data.item()[0]
W1_3_dps = W_234ML_xr_1deriv_sg.W_3ML_dps.data.item()[0]
W1_4_dps = W_234ML_xr_1deriv_sg.W_4ML_dps.data.item()[0]
'''
#S_234ML_xr_1deriv_sg
# pks N dps
S1_2_pks = S_234ML_xr_1deriv_sg.S_2ML_pks.data.item()[0]
S1_3_pks = S_234ML_xr_1deriv_sg.S_3ML_pks.data.item()[0]
S1_4_pks = S_234ML_xr_1deriv_sg.S_4ML_pks.data.item()[0]
S1_2_dps = S_234ML_xr_1deriv_sg.S_2ML_dps.data.item()[0]
S1_3_dps = S_234ML_xr_1deriv_sg.S_3ML_dps.data.item()[0]
S1_4_dps = S_234ML_xr_1deriv_sg.S_4ML_dps.data.item()[0]







# +
#W_234ML_xr_2deriv_sg
# pks N dps
'''
W2_2_pks = W_234ML_xr_2deriv_sg.W_2ML_pks.data.item()[0]
W2_3_pks = W_234ML_xr_2deriv_sg.W_3ML_pks.data.item()[0]
W2_4_pks = W_234ML_xr_2deriv_sg.W_4ML_pks.data.item()[0]
W2_2_dps = W_234ML_xr_2deriv_sg.W_2ML_dps.data.item()[0]
W2_3_dps = W_234ML_xr_2deriv_sg.W_3ML_dps.data.item()[0]
W2_4_dps = W_234ML_xr_2deriv_sg.W_4ML_dps.data.item()[0]
'''
#S_234ML_xr_1deriv_sg
# pks N dps
S2_2_pks = S_234ML_xr_2deriv_sg.S_2ML_pks.data.item()[0]
S2_3_pks = S_234ML_xr_2deriv_sg.S_3ML_pks.data.item()[0]
S2_4_pks = S_234ML_xr_2deriv_sg.S_4ML_pks.data.item()[0]
S2_2_dps = S_234ML_xr_2deriv_sg.S_2ML_dps.data.item()[0]
S2_3_dps = S_234ML_xr_2deriv_sg.S_3ML_dps.data.item()[0]
S2_4_dps = S_234ML_xr_2deriv_sg.S_4ML_dps.data.item()[0]






# -

S2_2_pks

# +
# filtering the peaks & dips 
# LDOS at peak position >  resolutin limit 
# set LDOS limit 
LIX_limit= 1E-14
# W type, 1st deriv peaks & dips
'''
W1_2ML_pks = W1_2_pks[(W_234ML_xr.W_2ML[W1_2_pks] > LIX_limit).values]
W1_3ML_pks = W1_3_pks[(W_234ML_xr.W_3ML[W1_3_pks] > LIX_limit).values]
W1_4ML_pks = W1_4_pks[(W_234ML_xr.W_4ML[W1_4_pks] > LIX_limit).values]

W1_2ML_dps = W1_2_dps[(W_234ML_xr.W_2ML[W1_2_dps] > LIX_limit).values]
W1_3ML_dps = W1_3_dps[(W_234ML_xr.W_3ML[W1_3_dps] > LIX_limit).values]
W1_4ML_dps = W1_4_dps[(W_234ML_xr.W_4ML[W1_4_dps] > LIX_limit).values]

# W type, 2nd deriv peaks & dips
W2_2ML_pks = W2_2_pks[(W_234ML_xr.W_2ML[W2_2_pks] > LIX_limit).values]
W2_3ML_pks = W2_3_pks[(W_234ML_xr.W_3ML[W2_3_pks] > LIX_limit).values]
W2_4ML_pks = W2_4_pks[(W_234ML_xr.W_4ML[W2_4_pks] > LIX_limit).values]

W2_2ML_dps = W2_2_dps[(W_234ML_xr.W_2ML[W2_2_dps] > LIX_limit).values]
W2_3ML_dps = W2_3_dps[(W_234ML_xr.W_3ML[W2_3_dps] > LIX_limit).values]
W2_4ML_dps = W2_4_dps[(W_234ML_xr.W_4ML[W2_4_dps] > LIX_limit).values]
'''
# S type, 1st deriv peaks & dips
S1_2ML_pks = S1_2_pks[(S_234ML_xr.S_2ML[S1_2_pks] > LIX_limit).values]
S1_3ML_pks = S1_3_pks[(S_234ML_xr.S_3ML[S1_3_pks] > LIX_limit).values]
S1_4ML_pks = S1_4_pks[(S_234ML_xr.S_4ML[S1_4_pks] > LIX_limit).values]

S1_2ML_dps = S1_2_dps[(S_234ML_xr.S_2ML[S1_2_dps] > LIX_limit).values]
S1_3ML_dps = S1_3_dps[(S_234ML_xr.S_3ML[S1_3_dps] > LIX_limit).values]
S1_4ML_dps = S1_4_dps[(S_234ML_xr.S_4ML[S1_4_dps] > LIX_limit).values]

# W type, 2nd deriv peaks & dips
S2_2ML_pks = S2_2_pks[(S_234ML_xr.S_2ML[S2_2_pks] > LIX_limit).values]
S2_3ML_pks = S2_3_pks[(S_234ML_xr.S_3ML[S2_3_pks] > LIX_limit).values]
S2_4ML_pks = S2_4_pks[(S_234ML_xr.S_4ML[S2_4_pks] > LIX_limit).values]

S2_2ML_dps = S2_2_dps[(S_234ML_xr.S_2ML[S2_2_dps] > LIX_limit).values]
S2_3ML_dps = S2_3_dps[(S_234ML_xr.S_3ML[S2_3_dps] > LIX_limit).values]
S2_4ML_dps = S2_4_dps[(S_234ML_xr.S_4ML[S2_4_dps] > LIX_limit).values]



#W1_2_pks
# -

# ## select peak positions 
#
# * delete the minor peak ( comes from ripples in the gap region)

# + tags=[]
'''
print(W2_2ML_dps)
fig,axs = plt.subplots()
axs2 =axs.twinx()


sns.lineplot(data = W_234ML_xr.to_dataframe(), x =  'bias_mV', y = 'W_2ML', ax=axs, color  = 'b')
sns.scatterplot(x =  W_234ML_xr.bias_mV[W2_2ML_dps_N], y = W_234ML_xr.W_2ML[W2_2ML_dps_N],ax=axs, color  = 'b')
axs.set_ylabel('LDOS', color='b')

sns.lineplot(data = W_234ML_xr_2deriv_sg.to_dataframe(), x =  'bias_mV', y = 'W_2ML', ax=axs2, color  = 'g')
sns.scatterplot(x = W_234ML_xr_2deriv_sg.bias_mV[W2_2ML_dps_N], y = W_234ML_xr_2deriv_sg.W_2ML[W2_2ML_dps_N],ax=axs2, color  = 'g')
axs2.set_ylabel('d2(LODS)/dV2', color='g')
W2_2ML_dps_N = W2_2ML_dps[[1,2,4,5]]
'''

# +

'''
print(W2_3ML_dps)
fig,axs = plt.subplots()
axs2 =axs.twinx()

sns.lineplot(data = W_234ML_xr.to_dataframe(), x =  'bias_mV', y = 'W_3ML', ax=axs, color  = 'b')
sns.scatterplot(x =  W_234ML_xr.bias_mV[W2_3ML_dps_N], y = W_234ML_xr.W_3ML[W2_3ML_dps_N],ax=axs, color  = 'b')
axs.set_ylabel('LDOS', color='b')

sns.lineplot(data = W_234ML_xr_2deriv_sg.to_dataframe(), x =  'bias_mV', y = 'W_3ML', ax=axs2, color  = 'g')
sns.scatterplot(x = W_234ML_xr_2deriv_sg.bias_mV[W2_3ML_dps_N], y = W_234ML_xr_2deriv_sg.W_3ML[W2_3ML_dps_N],ax=axs2, color  = 'g')
axs2.set_ylabel('d2(LODS)/dV2', color='g')
W2_3ML_dps_N = W2_3ML_dps[[2,3,5,6]]
'''
# -

'''
print(W2_4ML_dps)
fig,axs = plt.subplots()
axs2 =axs.twinx()


sns.lineplot(data = W_234ML_xr.to_dataframe(), x =  'bias_mV', y = 'W_4ML', ax=axs, color  = 'b')
sns.scatterplot(x =  W_234ML_xr.bias_mV[W2_4ML_dps_N], y = W_234ML_xr.W_4ML[W2_4ML_dps_N],ax=axs, color  = 'b')
axs.set_ylabel('LDOS', color='b')

sns.lineplot(data = W_234ML_xr_2deriv_sg.to_dataframe(), x =  'bias_mV', y = 'W_4ML', ax=axs2, color  = 'g')
sns.scatterplot(x = W_234ML_xr_2deriv_sg.bias_mV[W2_4ML_dps_N], y = W_234ML_xr_2deriv_sg.W_4ML[W2_4ML_dps_N],ax=axs2, color  = 'g')
axs2.set_ylabel('d2(LODS)/dV2', color='g')
W2_4ML_dps_N = W2_4ML_dps[[0,1,2,5,6]]
'''

# +
print(S2_2ML_dps)

# Check the S2_2ML_dps scatter  first 

# then use the 'S2_2ML_dps_N' scatter


fig,axs = plt.subplots()
axs2 =axs.twinx()


sns.lineplot(data = S_234ML_xr.to_dataframe(), x = 'bias_mV', y = 'S_2ML', ax=axs, color  = 'b')
sns.scatterplot(x =  S_234ML_xr.bias_mV[S2_2ML_dps_N], y = S_234ML_xr.S_2ML[S2_2ML_dps_N],ax=axs, color  = 'b')
axs.set_ylabel('LDOS', color='b')

sns.lineplot(data = S_234ML_xr_2deriv_sg.to_dataframe(), x =  'bias_mV', y = 'S_2ML', ax=axs2, color  = 'g')
sns.scatterplot(x = S_234ML_xr_2deriv_sg.bias_mV[S2_2ML_dps_N], y = S_234ML_xr_2deriv_sg.S_2ML[S2_2ML_dps_N], ax=axs2, color  = 'g')
axs2.set_ylabel('d2(LODS)/dV2', color='g')
S2_2ML_dps_N = S2_2ML_dps[[0,1,2,3,5,6]]

# +
print(S2_3ML_dps)
fig,axs = plt.subplots()
axs2 =axs.twinx()

sns.lineplot(data = S_234ML_xr.to_dataframe(), x =  'bias_mV', y = 'S_3ML', ax=axs, color  = 'b')
sns.scatterplot(x =  S_234ML_xr.bias_mV[S2_3ML_dps_N], y = S_234ML_xr.S_3ML[S2_3ML_dps_N],ax=axs, color  = 'b')
axs.set_ylabel('LDOS', color='b')

sns.lineplot(data = S_234ML_xr_2deriv_sg.to_dataframe(), x =  'bias_mV', y = 'S_3ML', ax=axs2, color  = 'g')
sns.scatterplot(x = S_234ML_xr_2deriv_sg.bias_mV[S2_3ML_dps_N], y = S_234ML_xr_2deriv_sg.S_3ML[S2_3ML_dps_N],ax=axs2, color  = 'g')
axs2.set_ylabel('d2(LODS)/dV2', color='g')
S2_3ML_dps_N = S2_3ML_dps[[0,1,2,3,4,5]]

# +
print(S2_4ML_dps)
fig,axs = plt.subplots()
axs2 =axs.twinx()


sns.lineplot(data = S_234ML_xr.to_dataframe(), x =  'bias_mV', y = 'S_4ML', ax=axs, color  = 'b')
sns.scatterplot(x =  S_234ML_xr.bias_mV[S2_4ML_dps_N], y = S_234ML_xr.S_4ML[S2_4ML_dps_N],ax=axs, color  = 'b')
axs.set_ylabel('LDOS', color='b')

sns.lineplot(data = S_234ML_xr_2deriv_sg.to_dataframe(), x =  'bias_mV', y = 'S_4ML', ax=axs2, color  = 'g')
sns.scatterplot(x = S_234ML_xr_2deriv_sg.bias_mV[S2_4ML_dps_N], y = S_234ML_xr_2deriv_sg.S_4ML[S2_4ML_dps_N],ax=axs2, color  = 'g')
axs2.set_ylabel('d2(LODS)/dV2', color='g')
S2_4ML_dps_N = S2_4ML_dps[[0,1,2,3,4,6,7,8]]
# -

# ## draw all 2,3,4 curves together 
# * with a offset 

# W stacks 
'''
fig,axs = plt.subplots()

offset_3ML = 1E-12
offset_4ML = 2E-12

sns.lineplot(data = W_234ML_xr.to_dataframe(), x =  'bias_mV', y = 'W_2ML', ax=axs, color  = 'b')
sns.scatterplot(x =  W_234ML_xr.bias_mV[W2_2ML_dps_N], y = W_234ML_xr.W_2ML[W2_2ML_dps_N],ax=axs, color  = 'b')
axs.set_ylabel('LDOS', color='b')

sns.lineplot(data = W_234ML_xr.to_dataframe()+offset_3ML, x =  'bias_mV', y = 'W_3ML', ax=axs, color  = 'b')
sns.scatterplot(x =  W_234ML_xr.bias_mV[W2_3ML_dps_N], y = W_234ML_xr.W_3ML[W2_3ML_dps_N]+offset_3ML,ax=axs, color  = 'b')
axs.set_ylabel('LDOS', color='b')

sns.lineplot(data = W_234ML_xr.to_dataframe()+offset_4ML, x =  'bias_mV', y = 'W_4ML', ax=axs, color  = 'b')
sns.scatterplot(x =  W_234ML_xr.bias_mV[W2_4ML_dps_N], y = W_234ML_xr.W_4ML[W2_4ML_dps_N]+offset_4ML,ax=axs, color  = 'b')
axs.set_ylabel('LDOS', color='b')
'''

# +
# S  stacks 
fig,axs = plt.subplots()

offset_3ML = 2E-12
offset_4ML = 4E-12

sns.lineplot(data = S_234ML_xr.to_dataframe(), x = 'bias_mV', y = 'S_2ML', ax=axs, color  = 'k')
sns.scatterplot(x =  S_234ML_xr.bias_mV[S2_2ML_dps_N], y = S_234ML_xr.S_2ML[S2_2ML_dps_N],ax=axs, color  = 'k')
axs.set_ylabel('LDOS', color='k')
sns.lineplot(data = S_234ML_xr.to_dataframe()+offset_3ML, x =  'bias_mV', y = 'S_3ML', ax=axs, color  = 'k')
sns.scatterplot(x =  S_234ML_xr.bias_mV[S2_3ML_dps_N], y = S_234ML_xr.S_3ML[S2_3ML_dps_N]+offset_3ML,ax=axs, color  = 'k')
axs.set_ylabel('LDOS', color='k')
sns.lineplot(data = S_234ML_xr.to_dataframe()+offset_4ML, x =  'bias_mV', y = 'S_4ML', ax=axs, color  = 'k')
sns.scatterplot(x =  S_234ML_xr.bias_mV[S2_4ML_dps_N], y = S_234ML_xr.S_4ML[S2_4ML_dps_N]+offset_4ML,ax=axs, color  = 'k')
axs.set_ylabel('LDOS', color='k')

## 3ML --> wrinkle effect 

# -













# + [markdown] tags=[]
# # 2. ROI (Aera)  selection 
# ## 2-1 Bound Box
#     * using Bounding Box 
#     * live 
# ## 2-2 Lasso Selection 
# * using selection1D 
#     * using xr for groupby  & merge (join = inner)
#     use the "combine_by_coords" in XR

# + [markdown] tags=[]
#
# ## 2-1-1 bokeh plot &  Bound box selection + get point selection 
#

# +
hv.extension('bokeh')

grid_channel_hv = hv.Dataset(grid_3D.I_fb)

# bias_mV slicing
dmap_plane  =  [ "X","Y"]
dmap = grid_channel_hv.to(hv.Image,
                          kdims = dmap_plane,
                          dynamic = True )
dmap.opts(colorbar = True,
          cmap = 'bwr',
          frame_width = 400,
          aspect = 1 )#.relabel('XY plane slicing: ')

grid_channel_hv_image  = hv.Dataset(grid_3D.I_fb.isel(bias_mV = 0)).relabel('for BBox selection : ')

bbox_points = hv.Points(grid_channel_hv_image).opts(frame_width = 400,
                                                    color = 'k',
                                                    aspect = 1,
                                                    alpha = 0.1,                                   
                                                    tools=['box_select'])

bound_box = hv.streams.BoundsXY(source = bbox_points,
                                bounds=(0,0,0,0))
dmap*bbox_points

## hv.DynamicMap( 뒤에는 function 이 와야함), streams  로 해당 영역을 지정.( or 함수의 입력정보 지정) 
# averaged curve 를 그리기 위해서 해당영역을  xr  에서 average  해야함.. 

# curve 의 area 로 error bar도 같이 그릴것.. 



# + [markdown] tags=[]
# ## 2-1-2 averaged STS of the selected Bound Box region

# +
#bound_box.bounds
x_bounds_msk = (grid_3D.X > bound_box.bounds[0] ) & (grid_3D.X < bound_box.bounds[2])
y_bounds_msk = (grid_3D.Y > bound_box.bounds[1] ) & (grid_3D.Y < bound_box.bounds[3])
grid_3D_bbox = grid_3D.where(grid_3D.X[x_bounds_msk] + grid_3D.Y[y_bounds_msk])


dmap* hv.Dataset(grid_3D_bbox).to(hv.Image, kdims = ['X','Y'],
                                  dynamic = True).opts(cmap = 'cwr')\
+ hv.Dataset(grid_3D_bbox.LIX_fb.mean(dim = ['X','Y'])).to(hv.Curve)

#sns_curve = sns.lineplot(x = "bias_mV",
#                          y = "LIX_fb", 
#                           data = grid_3D_bbox.to_dataframe())
# at present, for the confidence interval ==> use the sns

# + [markdown] tags=[]
# ## 2-2-1 bokeh plot & Lasso Selection + get point_lists

# +
hv.extension('bokeh')


grid_channel_hv = hv.Dataset(grid_3D.I_fb)

# bias_mV slicing
dmap_plane  = ["X","Y"]
dmap = grid_channel_hv.to(hv.Image,
                          kdims = dmap_plane,
                          dynamic = True )
dmap.opts(colorbar = True,
          cmap = 'bwr',
          frame_width = 400,
          aspect = 1).relabel('XY plane slicing: ')



grid_channel_hv_image = hv.Dataset(grid_3D.I_fb.isel(bias_mV = 0))

grid_channel_hv_points = hv.Points(grid_channel_hv_image).opts(frame_width = 400,  
                                   aspect = 1, alpha = 0.1,                                   
                                   tools=['box_select', 'lasso_select']
                                  )

slct_pts = hv.streams.Selection1D(source=grid_channel_hv_points)

dmap*grid_channel_hv_image*grid_channel_hv_points

# +
#slct_pts
pts = grid_channel_hv_points.iloc[slct_pts.index].dframe().set_index(['X', 'Y'])

pts_xr = xr.Dataset.from_dataframe(pts)

grid_3D_slct_pts = xr.combine_by_coords ([grid_3D, pts_xr], compat = 'override', join = 'inner')
#y_pts = points.iloc[slct_pts.index].dframe().Y
#grid_3D.sel(X = x_pts,Y = y_pts)
#grid_3D.I_fb.isel(bias_mV = 0).plot()
#grid_3D_slct_pts

dmap* hv.Dataset(grid_3D_slct_pts).to(hv.Image, kdims = ['X','Y'],
                                  dynamic = True).opts(cmap = 'cwr')\
+ hv.Dataset(grid_3D_slct_pts.LIX_fb.mean(dim = ['X','Y'])).to(hv.Curve)

# + [markdown] tags=[]
# # At present, Interactive plot for 3D data map is not working 
#     * For interactive XY positioning & LDOS curve $\to$ use the 'streamz'
# ## instead, use the two-step selection
#     
#
# ## use the holoview to select XY position 
# * select XY with slider 
#
#

# +
#grid_3D.X.values.min()
from panel.interact import interact
sliderX = pnw.IntSlider(name='X', 
                       start = 0 ,
                       end = grid_3D.X.shape[0]) 
sliderY = pnw.IntSlider(name='Y', 
                       start = 0 ,
                       end = grid_3D.Y.shape[0]) 
# 여기서  control 하면 아래쪽에서 sliderX 혹은 slider Y를 불러내도 똑같이 반응함. \

sliderX_v_intact = interact(lambda x:  grid_3D.X[x].values, x =sliderX)[1]
sliderY_v_intact = interact(lambda y:  grid_3D.Y[y].values, y =sliderY)[1]
# 이렇게 하면 값들을 확인 할 수 있으나, 
# interactive 하게 결정된 값이 pn.Row type object로   integer 변환이 되지 않는다. 
# interactive 하게 변하는 값이 아닌 한번 결정된값을 불러올때는 단순히 value  를 붙여서 꺼낼수 있지만, 
# interacive 한 stream을 연결시키기위해서는 다른 방법이 필요함. 
pn.Column(interact(lambda x:  grid_3D.X[x].values, x =sliderX), interact(lambda y: grid_3D.Y[y].values, y =sliderY))
# Do not exceed the max Limit ==> error

# + [markdown] tags=[]
# ## STS curve at XY point 
# -

# Fixed value from slider XY
sliderX_v = grid_3D.X[sliderX.value].values
sliderY_v = grid_3D.Y[sliderY.value].values
# add HLine & VLines # H line is for Y & V line is for X!! 
dmap_slideXY = dmap *hv.VLine(sliderX_v).opts(
    color = 'green', line_width = 4, alpha =0.5)*hv.HLine(
    sliderY_v).opts(
    color = 'yellow', line_width = 4, alpha =0.5)
grid_3D_selXY= grid_3D.isel(X = sliderX.value, Y= sliderY.value)
grid_3D_I_plot = hv.Curve(grid_3D_selXY.I_fb).relabel('I_fb')
grid_3D_LIX_plot = hv.Curve(grid_3D_selXY.LIX_fb).relabel('LIX_fb')
pn.Column(dmap_slideXY, grid_3D_I_plot +grid_3D_LIX_plot)

# ## Build a function for plotting semi-interactive plot 
# * choose the XY points 
# * choose the channels 
# * choose the smoothing level 
#
# ### Plot map & curves 
# * dmap * HLine * VLline + I Curve or LIX Curve 
# ### Plot line profiles (H&V) 
# * dmap + Line profile along the HLine & VLine 
# * If I I need a tilted line profile -> use the rotation 
#
# ### To Be Added
#

# + [markdown] tags=[] jp-MarkdownHeadingCollapsed=true
# # check the  Horizontal & Vertical Line Profiles 
# -

pn.Column(interact(lambda x:  grid_3D.X[x].values, x =sliderX), interact(lambda y: grid_3D.Y[y].values, y =sliderY))


# +


# Fixed value from slider XY
sliderX_v = grid_3D.X[sliderX.value].values
sliderY_v = grid_3D.Y[sliderY.value].values
##sliderX_v, sliderY


data_channel  = 'LIX_fb'
#Choose the plot channel line profile test 


grid_3D_Hline_profile = grid_3D.isel(Y = sliderY.value)[data_channel]
grid_3D_Vline_profile = grid_3D.isel(X = sliderX.value)[data_channel]

dmap_hline = hv.Dataset(grid_3D_Hline_profile).to(hv.Curve,
                          kdims = ['X'],
                          dynamic = True ).relabel('Horizontal Line Profile')
dmap_vline = hv.Dataset(grid_3D_Vline_profile).to(hv.Curve,
                          kdims = ['Y'],
                          dynamic = True ).relabel('Vertical Line Profile')
# get line profile from data set 

hv.extension('bokeh')
grid_channel_hv = hv.Dataset(grid_3D[data_channel])
# choose the 'LIX_fb'  for plotting 
# bias_mV slicing
dmap_plane  = ["X","Y"]
dmap = grid_channel_hv.to(hv.Image,
                          kdims = dmap_plane,
                          dynamic = True ).relabel('XY plane slicing: '+', channel : '+data_channel)

dmap.opts(colorbar = True,
          cmap = 'bwr',
          frame_width = 400,
          aspect = 1)
dmap_slideXY = dmap*hv.VLine(sliderX_v).opts(
    color = 'green', line_width = 4, alpha =0.5)*hv.HLine(
    sliderY_v).opts(
    color = 'yellow', line_width = 4, alpha =0.5)


dmap_slideXY+ dmap_hline+dmap_vline



# +
#grid_3D.X.values.max()
#help(pnw.DiscreteSlider)
data_varis_list = list(grid_3D_gap.variables.mapping.keys())
data_varis_list



 
# -


# * current issue : how to plot XY plane instead of YX plane 
#

# + [markdown] tags=[]
# # Signal Treatments 
# ## Savatzky-Golay smoothig 
#     * use the list comprehension for the sg-smoothing 

# +
grid_3D_sg = savgolFilter_xr(grid_3D, window_length = 7, polyorder = 3)
## savatzk-golay filtered 

# test the curve 
#grid_3D.LIX_unit_calc.isel(X = 20, Y = 20).plot()
#grid_3D_sg.LIX_unit_calc.isel(X = 20, Y = 20).plot()

# + [markdown] tags=[]
# ## Numerical derivative 
#     * use the xarray APIs 
# ### Compare the sg filtered result & derivative result 
#     
# -

# grid_3D -> sg -> derivative 
grid_3D_sg = savgolFilter_xr(grid_3D, window_length = 21, polyorder = 3)
grid_3D_sg_1deriv = grid_3D_sg.differentiate('bias_mV')
grid_3D_sg_1deriv_sg = savgolFilter_xr(grid_3D_sg_1deriv, window_length = 21, polyorder = 3)
grid_3D_sg_2deriv = grid_3D_sg_1deriv_sg.differentiate('bias_mV')

#pn.Row(sliderX,sliderY)
pn.Row(interact(lambda x:  grid_3D.X[x].values, x =sliderX), interact(lambda y: grid_3D.Y[y].values, y =sliderY))

grid_3D

# +
# Fixed value from slider XY
sliderX_v = grid_3D.X[sliderX.value].values
sliderY_v = grid_3D.Y[sliderY.value].values


LIX_unit_calc_curve = hv.Curve(grid_3D_sg.isel(X = sliderX.value).LIX_fb.values[0]).opts(axiswise=True, ylabel='LDOS (A/V)', title = 'LDOS')
d_LIX_unit_calc_curve_dV =  hv.Curve(grid_3D_sg_1deriv.isel(X = sliderX.value).LIX_fb[0]).opts(axiswise=True, ylabel='LDOS (A/V)', title = 'd(LDOS)/dV')
d2_LIX_unit_calc_curve_dV2 = hv.Curve(grid_3D_sg_2deriv.isel(X = sliderX.value).LIX_fb[0]).opts(axiswise=True, ylabel='LDOS (A/V)', title = 'd2(LDOS)/dV2')
LIX_unit_calc_curve + d_LIX_unit_calc_curve_dV + d2_LIX_unit_calc_curve_dV2



# + [markdown] tags=[]
# ## peak finding (scipy)
#     * use the list comprehension for peak finding 
# * function __find_peaks_xr__
# -

grid_3D_sg_pks = find_peaks_xr(grid_3D_sg)
#grid_3D_sg_pks.LIX_fb_peaks

grid_3D_sg_pks

# + [markdown] tags=[] jp-MarkdownHeadingCollapsed=true tags=[]
# ## Handing peaks with different numbers 
# * use the np.pad & count the max number of peak
# * fill the np.nan the empty 
# * function : __peak_pad__
#     
#    
# -

grid_3D_sg_pks_pad = peak_pad(grid_3D_sg_pks)
grid_3D_sg_pks_pad

# + [markdown] tags=[] jp-MarkdownHeadingCollapsed=true tags=[]
# ## counting peaks for after derivatives 
#
# * original data : grid_3D
# * find gap data : grid_3D_gap
# * smoothing data : grid_3D_gap_sg
# * derivative data : grid_3D_gap_sg_1deriv, grid_3D_gap_sg_2deriv
# * find_peaks  : grid_3D_sg_1deriv_sg_pks
# * find_peaks& padding  : grid_3D_sg_1deriv_sg_pks_pad
#
#
# ### LDOS in data 
# * grid_3D_gap_sg
# ### for the peaks in LDOS 
# * grid_3D_sg_pks_pad
# ###  d(LDOS)/dV IETS 
# * grid_3D_sg_1deriv_sg
# * peaks & dips in IETS : d2I/dV2 values 
#     * grid_3D_sg_1deriv_sg_pks_pad
#     * grid_3D_sg_1deriv_sg_dps_pad
# ### To find LDOS peaks $\to$ 2nd derivative LDOS  dips! (for acurate detection)
# * grid_3D_sg_2deriv_sg
#     * grid_3D_sg_2deriv_sg_pks_pad
#     * grid_3D_sg_2deriv_sg_dps_pad

# +
# grid_3D -> sg -> derivative 
grid_3D_sg = savgolFilter_xr(grid_3D, window_length = 31, polyorder = 3)
grid_3D_sg_1deriv = grid_3D_sg.differentiate('bias_mV')
grid_3D_sg_1deriv_sg = savgolFilter_xr(grid_3D_sg_1deriv, window_length = 31, polyorder = 3)
grid_3D_sg_2deriv = grid_3D_sg_1deriv_sg.differentiate('bias_mV')

# d(LDOS)dV pks & dps
grid_3D_sg_1deriv_sg_pks = find_peaks_xr(grid_3D_sg_1deriv_sg)
grid_3D_sg_1deriv_sg_pks_pad = peak_pad(grid_3D_sg_1deriv_sg_pks)

grid_3D_sg_1deriv_sg_dps = find_peaks_xr(-1* grid_3D_sg_1deriv_sg)
grid_3D_sg_1deriv_sg_dps_pad = peak_pad(grid_3D_sg_1deriv_sg_dps)

# d2(LDOS)dV2 pks & dps

grid_3D_sg_2deriv_sg =  savgolFilter_xr(grid_3D_sg_2deriv)
grid_3D_sg_2deriv_sg_pks = find_peaks_xr(grid_3D_sg_2deriv_sg)
grid_3D_sg_2deriv_sg_pks_pad = peak_pad(grid_3D_sg_2deriv_sg_pks)

grid_3D_sg_2deriv_sg_dps = find_peaks_xr(-1* grid_3D_sg_2deriv_sg)
grid_3D_sg_2deriv_sg_dps_pad = peak_pad(grid_3D_sg_2deriv_sg_dps)


# + [markdown] tags=[]
# ## recall XY points 
# -

# # After Treatments 
# * **grid_3D** : original 
# * **grid_3D_sg** : after savatzky golray filter (smoothing) 
# * **grid_3D_sg_1deriv** : after 1st derivative 
#     * I $\to$ dI/dV
#     * LIX $\to$ d(LIX)/dV
#     * dIdV $\to$ d2I/dV2
#     * LDOS  $\to$ d2I/dV2
# * **grid_3D_sg_2deriv_sg** : after 2nd derivative 
#     * I $\to$ d2I/dV2
#     * LIX $\to$ d2(LIX)/dV2
#     * dIdV $\to$ d3I/dV3 ( for dI/dV peak detection)
#     * LDOS  $\to$ d2I/dV2  ( for d(LDOS)/dV peak detection)
#

# +
#pn.Row(interact(lambda x:  grid_3D.X[x].values, x =sliderX), interact(lambda y: grid_3D.Y[y].values, y =sliderY))
# Fixed value from slider XY

pn.Row(sliderX,sliderY)

#grid_3D_sg_2deriv_sg.LIX_unit_calc


# +


# Fixed value from slider XY
sliderX_v = grid_3D.X[sliderX.value].values
sliderY_v = grid_3D.Y[sliderY.value].values
##sliderX_v, sliderY


data_channel  = 'LIX_fb'
#Choose the plot channel line profile test 


grid_3D_Hline_profile = grid_3D_sg_2deriv_sg.isel(Y = sliderY.value)[data_channel]
grid_3D_Vline_profile = grid_3D_sg_2deriv_sg.isel(X = sliderX.value)[data_channel]

dmap_hline = hv.Dataset(grid_3D_Hline_profile).to(hv.Curve,
                          kdims = ['X'],
                          dynamic = True ).relabel('Horizontal Line Profile')
dmap_vline = hv.Dataset(grid_3D_Vline_profile).to(hv.Curve,
                          kdims = ['Y'],
                          dynamic = True ).relabel('Vertical Line Profile')
# get line profile from data set 

hv.extension('bokeh')
grid_channel_hv = hv.Dataset(grid_3D_sg_2deriv_sg[data_channel])
# choose the 'LIX_fb'  for plotting 
# bias_mV slicing
dmap_plane  = ["X","Y"]
dmap = grid_channel_hv.to(hv.Image,
                          kdims = dmap_plane,
                          dynamic = True ).relabel('XY plane slicing: '+', channel : '+data_channel)

dmap.opts(colorbar = True,
          cmap = 'bwr',
          frame_width = 400,
          aspect = 1)
dmap_slideXY = dmap*hv.VLine(sliderX_v).opts(
    color = 'green', line_width = 4, alpha =0.5)*hv.HLine(
    sliderY_v).opts(
    color = 'yellow', line_width = 4, alpha =0.5)


dmap_slideXY+ dmap_hline+dmap_vline

# -


z_profile  = (PlaneFitY_xr(grid_topo) - PlaneFitY_xr(grid_topo).min())
z_profile.topography.plot()

# + [markdown] tags=[]
# #### 2ML peak detection 

# +
STS_P_2ML_1deriv = STS_P_2ML.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()
STS_P_2ML_2deriv = STS_P_2ML_1deriv.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()


STS_P_2ML_2deriv_dps,_ =  sp.signal.find_peaks(-1*STS_P_2ML_2deriv,prominence=2E-18)
# use the prominence to choose the proper peak position 
STS_P_2ML_2deriv_dps_width = sp.signal.peak_widths(STS_P_2ML_2deriv, STS_P_2ML_2deriv_dps, rel_height=0.5)

STS_P_2ML_2deriv_dps_filtered = STS_P_2ML_2deriv_dps[
    (STS_P_2ML_2deriv.bias_mV[STS_P_2ML_2deriv_dps].values
     > grid_3D_gap_STS.groupby('layer')['2ML'].mean(dim='stacked_X_Y').CBM_I_mV.values )
    +(STS_P_2ML_2deriv.bias_mV[STS_P_2ML_2deriv_dps].values <
      grid_3D_gap_STS.groupby('layer')['2ML'].mean(dim='stacked_X_Y').VBM_I_mV.values)]

# dip positions in 2nd derivative 
STS_P_2ML_2deriv_dps_plt = hv.Curve(
    hv.Dataset(STS_P_2ML_2deriv)).relabel('d2(LDOS)/dV2').opts(
    axiswise = True)*hv.Scatter(
    hv.Dataset(STS_P_2ML_2deriv[STS_P_2ML_2deriv_dps])).opts(
    color= 'k', marker="cross",size = 10).opts( 
    axiswise = True)
# dip positions apply to LDOS 
STS_P_2ML_pks_plt = hv.Curve(
    hv.Dataset(STS_P_2ML)).opts(axiswise = True).relabel('LDOS')*hv.Scatter(
    hv.Dataset(STS_P_2ML[STS_P_2ML_2deriv_dps_filtered])).opts(
    color= 'k', marker="cross",size = 10, axiswise = True).relabel('LDOS')

STS_P_2ML_2deriv_dps_plt.relabel('d2(LDOS)/dV2').opts(axiswise = True) + STS_P_2ML_pks_plt.opts(axiswise = True)

# + [markdown] tags=[]
# #### 3ML peak detection 
# -

grid_3D_gap_STS.groupby('layer')['3ML'].mean(dim='stacked_X_Y').VBM_LIX_mV.values

# + jp-MarkdownHeadingCollapsed=true tags=[]
STS_P_3ML_1deriv = STS_P_3ML.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()
STS_P_3ML_2deriv = STS_P_3ML_1deriv.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()


STS_P_3ML_2deriv_dps,_ =  sp.signal.find_peaks(-1*STS_P_3ML_2deriv,prominence=2E-18)
# use the prominence to choose the proper peak position 
STS_P_3ML_2deriv_dps_width = sp.signal.peak_widths(STS_P_3ML_2deriv, STS_P_3ML_2deriv_dps, rel_height=0.5)

STS_P_3ML_2deriv_dps_filtered = STS_P_3ML_2deriv_dps[
    (STS_P_3ML_2deriv.bias_mV[STS_P_3ML_2deriv_dps].values
     > grid_3D_gap_STS.groupby('layer')['3ML'].mean(dim='stacked_X_Y').CBM_LIX_mV.values )
    +(STS_P_3ML_2deriv.bias_mV[STS_P_3ML_2deriv_dps].values <
      grid_3D_gap_STS.groupby('layer')['3ML'].mean(dim='stacked_X_Y').VBM_LIX_mV.values)]

# dip positions in 2nd derivative 
STS_P_3ML_2deriv_dps_plt = hv.Curve(
    hv.Dataset(STS_P_3ML_2deriv)).relabel('d2(LDOS)/dV2').opts(
    axiswise = True)*hv.Scatter(
    hv.Dataset(STS_P_3ML_2deriv[STS_P_3ML_2deriv_dps_filtered])).opts(
    color= 'k', marker="cross",size = 10).opts( 
    axiswise = True)
# dip positions apply to LDOS 
STS_P_3ML_pks_plt = hv.Curve(
    hv.Dataset(STS_P_3ML)).opts(axiswise = True).relabel('LDOS')*hv.Scatter(
    hv.Dataset(STS_P_3ML[STS_P_3ML_2deriv_dps_filtered])).opts(
    color= 'k', marker="cross",size = 10, axiswise = True).relabel('LDOS')

STS_P_3ML_2deriv_dps_plt.relabel('d2(LDOS)/dV2').opts(axiswise = True) + STS_P_3ML_pks_plt.opts(axiswise = True)
# -

# #### 4ML peak detection

# +
STS_P_4ML_1deriv = STS_P_4ML.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()
STS_P_4ML_2deriv = STS_P_4ML_1deriv.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()


STS_P_4ML_2deriv_dps,_ =  sp.signal.find_peaks(-1*STS_P_4ML_2deriv,prominence=2E-18)
# use the prominence to choose the proper peak position 
STS_P_4ML_2deriv_dps_width = sp.signal.peak_widths(STS_P_4ML_2deriv, STS_P_4ML_2deriv_dps, rel_height=0.5)

STS_P_4ML_2deriv_dps_filtered = STS_P_4ML_2deriv_dps[
    (STS_P_4ML_2deriv.bias_mV[STS_P_4ML_2deriv_dps].values
     > grid_3D_gap_STS.groupby('layer')['4ML'].mean(dim='stacked_X_Y').CBM_LIX_mV.values )
    +(STS_P_4ML_2deriv.bias_mV[STS_P_4ML_2deriv_dps].values <
      grid_3D_gap_STS.groupby('layer')['4ML'].mean(dim='stacked_X_Y').VBM_LIX_mV.values)]

# dip positions in 2nd derivative 
STS_P_4ML_2deriv_dps_plt = hv.Curve(
    hv.Dataset(STS_P_4ML_2deriv)).relabel('d2(LDOS)/dV2').opts(
    axiswise = True)*hv.Scatter(
    hv.Dataset(STS_P_4ML_2deriv[STS_P_4ML_2deriv_dps])).opts(
    color= 'k', marker="cross",size = 10).opts( 
    axiswise = True)
# dip positions apply to LDOS 
STS_P_4ML_pks_plt = hv.Curve(
    hv.Dataset(STS_P_4ML)).opts(axiswise = True).relabel('LDOS')*hv.Scatter(
    hv.Dataset(STS_P_4ML[STS_P_4ML_2deriv_dps_filtered])).opts(
    color= 'k', marker="cross",size = 10, axiswise = True).relabel('LDOS')


STS_P_4ML_2deriv_dps_plt.relabel('d2(LDOS)/dV2').opts(axiswise = True) + STS_P_4ML_pks_plt.opts(axiswise = True)

# +


STS_P_2ML_pks_plt = hv.Curve(
    hv.Dataset(STS_P_2ML), label = '2ML').opts(axiswise = True)*hv.Scatter(
    hv.Dataset(STS_P_2ML[STS_P_2ML_2deriv_dps_filtered])).opts(
     marker="cross",size = 10, axiswise = True)


STS_P_3ML_pks_plt = hv.Curve(
    hv.Dataset(STS_P_3ML), label = '3ML').opts(axiswise = True)*hv.Scatter(
    hv.Dataset(STS_P_3ML[STS_P_3ML_2deriv_dps_filtered])).opts(
     marker="cross",size = 10, axiswise = True)

STS_P_4ML_pks_plt = hv.Curve(
    hv.Dataset(STS_P_4ML), label = '4ML').opts(axiswise = True)*hv.Scatter(
    hv.Dataset(STS_P_4ML[STS_P_4ML_2deriv_dps_filtered])).opts(
     marker="cross",size = 10, axiswise = True)


(STS_P_2ML_pks_plt*STS_P_3ML_pks_plt*STS_P_4ML_pks_plt)
# -

## check the peaks with histogram 
grid_3D_gap_STS.groupby('layer')['3ML'].to_dataframe().LDOS_fb.to_xarray()

# +
grid_3D_gap_sg_2deriv_sg_dps_pad.layer.values = grid_3D_gap.layer.copy(deep = True)
test_2ML = grid_3D_gap_sg_2deriv_sg_dps_pad.groupby('layer')['2ML'].unstack().LDOS_fb_peaks_pad
test_3ML = grid_3D_gap_sg_2deriv_sg_dps_pad.groupby('layer')['3ML'].unstack().LDOS_fb_peaks_pad
test_4ML = grid_3D_gap_sg_2deriv_sg_dps_pad.groupby('layer')['4ML'].unstack().LDOS_fb_peaks_pad

test_2ML_hv  = hv.Dataset(test_2ML)
test_3ML_hv  = hv.Dataset(test_3ML)
test_4ML_hv  = hv.Dataset(test_4ML)

test_4ML_hv.to(hv.Image, kdims= ['X','Y'], dynamic = True)+test_3ML_hv.to(hv.Image, kdims= ['X','Y'], dynamic = True)+test_2ML_hv.to(hv.Image, kdims= ['X','Y'], dynamic = True)

# -

grid_3D_gap_sg_2deriv_sg_dps_pad

# # Use the Lasso & select area based on dynamic map 

# + [markdown] tags=[]
# ## 2-2-1 bokeh plot & Lasso Selection + get point_lists

# +
hv.extension('bokeh')


grid_channel_hv = hv.Dataset(grid_3D_gap_sg_2deriv_sg_dps_pad.LDOS_fb)

# bias_mV slicing
dmap_plane  = ["X","Y"]
dmap = grid_channel_hv.to(hv.Image,
                          kdims = dmap_plane,
                          dynamic = True )
dmap.opts(colorbar = True,
          cmap = 'bwr',
          frame_width = 200,
          aspect =0.5 ).relabel('XY plane slicing: ')



grid_channel_hv_image = hv.Dataset(grid_3D.I_fb.isel(bias_mV = 0))

grid_channel_hv_points = hv.Points(grid_channel_hv_image).opts(frame_width = 200,
                                   aspect =0.5, alpha = 0.1,                                   
                                   tools=['box_select', 'lasso_select']
                                  )

slct_pts = hv.streams.Selection1D(source=grid_channel_hv_points)

dmap*grid_channel_hv_image*grid_channel_hv_points
# -

#slct_pts
### 
"""
pts_2ML1 = grid_channel_hv_points.iloc[slct_pts.index].dframe().set_index(['X', 'Y'])
pts_2ML1_xr = xr.Dataset.from_dataframe(pts_2ML1)
grid_3D_slct_pts_2ML1 = xr.combine_by_coords ([grid_3D, pts_2ML1_xr], compat = 'override', join = 'inner')

pts_3ML1 = grid_channel_hv_points.iloc[slct_pts.index].dframe().set_index(['X', 'Y'])
pts_3ML1_xr = xr.Dataset.from_dataframe(pts_3ML1)
grid_3D_slct_pts_3ML1 = xr.combine_by_coords ([grid_3D, pts_3ML1_xr], compat = 'override', join = 'inner')
"""
pts_4ML1 = grid_channel_hv_points.iloc[slct_pts.index].dframe().set_index(['X', 'Y'])
pts_4ML1_xr = xr.Dataset.from_dataframe(pts_4ML1)
grid_3D_slct_pts_4ML1 = xr.combine_by_coords ([grid_3D, pts_4ML1_xr], compat = 'override', join = 'inner')


hv.Dataset(grid_3D_slct_pts_2ML1.LIX_fb.mean(dim = ['X','Y'])).to(hv.Curve)*hv.Dataset(grid_3D_slct_pts_3ML1.LIX_fb.mean(dim = ['X','Y'])).to(hv.Curve)*hv.Dataset(grid_3D_slct_pts_4ML1.LIX_fb.mean(dim = ['X','Y'])).to(hv.Curve)

# +
#y_pts = points.iloc[slct_pts.index].dframe().Y
#grid_3D.sel(X = x_pts,Y = y_pts)
#grid_3D.I_fb.isel(bias_mV = 0).plot()
#grid_3D_slct_pts

dmap* hv.Dataset(grid_3D_slct_pts).to(hv.Image, kdims = ['X','Y'],
                                  dynamic = True).opts(cmap = 'cwr')\
+ hv.Dataset(grid_3D_slct_pts.LIX_fb.mean(dim = ['X','Y'])).to(hv.Curve)
# -

pts_3ML1

# +
#.LDOS_fb_VB_peaks_pad.to_dataframe().melt()

test[test.value > 10]


sns.histplot(data=test[test.value > 10], x="value", kde=True, bins = 200)


#grid_3D_gap_sg_2deriv_sg_dps_pad.LDOS_fb_peaks_pad.to_dataframe().hist(bins = 200)

# + [markdown] tags=[]
# #### Chosse the point spectra using mask 

# +
#lineSTS_P.iloc[60:,:].shape
#lineSTS_P.iloc[60:,:]
#idx = pd.IndexSlice 
mask_terrace  = (lineSTS_P.layer !='boundary') 
mask_X  = (lineSTS_P.X < 0.7E-7) 
lineSTS_P[mask_terrace &mask_X]

grid_3D_gap_STS = xr.merge([grid_3D_gap_sg[['I_fb','LDOS_fb','CBM_I_mV','VBM_I_mV','CBM_LIX_mV','VBM_LIX_mV','layer']] , grid_3D_gap_sg_2deriv_sg_dps_pad[['LDOS_fb_peaks']]]).dropna(dim = 'X')

# -

sns.lineplot(data  = lineSTS_P, x = 'bias_mV', y='LDOS_fb', hue = 'layer')

# ## choose the terrace regions only 

#z_profile[z_profile.edge == 0].topography.plot()
z_profile_plot = sns.lineplot(x = "X", y = "topography", data = z_profile)
z_profile_plot_terrace = sns.scatterplot(x = "X", y = "topography", data = z_profile[z_profile.edge == 0], color= 'k')


# #### choose the uphill terraces based on topography 

z_profile_up = z_profile.reset_index()[z_profile.reset_index().X < 0.7E-7].set_index(['X','Y'])
z_profile_up_terrace = z_profile_up[z_profile_up.edge == 0]
z_profile_up_terrace
z_profile_up_plot = sns.lineplot(x = "X", y = "topography", data = z_profile_up)
z_profile_up_plot_terrace = sns.scatterplot(x = "X", y = "topography", data = z_profile_up_terrace, color= 'k')



# +
layer_N = z_profile_up_terrace.to_xarray() 
# pd.qcut 활용한 결과를 xarray로 바꿈. 

grid_3D_gap_sg['layer'] = layer_N.label
#grid_3D_sg

lineSTS_P = grid_3D_gap_sg[['I_fb','LDOS_fb','layer']].to_dataframe()
#lineSTS_P

sns.lineplot(data  = lineSTS_P, x = 'bias_mV', y='LDOS_fb', hue = 'layer')
# -

# ### select curves
# ### find peaks 
#
#
#

#target_curve = peak_detection_curve[data_ch]
#target_curve.bias_mV
target_curve_HOPG = lineSTS_P[lineSTS_P.layer == 'HOPG'].LDOS_fb.groupby('bias_mV').mean()
target_curve_1ML = lineSTS_P[lineSTS_P.layer == '1ML'].LDOS_fb.groupby('bias_mV').mean()
target_curve_HOPG.plot()
target_curve_1ML.plot()

# +

target_curve_HOPG = target_curve_HOPG* 1E9
target_curve_HOPG.index  = target_curve_HOPG.index*1E-3
target_curve_HOPG.index.names  = ['Bias (V)']
target_curve_HOPG

target_curve_1ML = target_curve_1ML* 1E9
target_curve_1ML.index  = target_curve_1ML.index*1E-3
target_curve_1ML.index.names  = ['Bias (V)']
target_curve_1ML
# -

# hv plot 


hv.extension('matplotlib')
hv.Curve(target_curve_HOPG).opts(
    color = 'black')* hv.Curve(
    target_curve_1ML).opts(
    color = 'blue').opts(
    logy= True,
    ylabel  = 'dI/dV (nA/V)',
    xlabel = 'Bias (V)',
    fontsize= 'x-large', 
    fig_size= 200)

# ### HOPG  & 1ML peaks 
# #### find peaks unig sp.find_peaks

# +
peaks,_ =  sp.signal.find_peaks(target_curve_HOPG,height = 1E-5, prominence =  1E-6)
pk_width = sp.signal.peak_widths(target_curve_HOPG, peaks, rel_height=0.5)

hv.extension('bokeh')
#hv.extension('matplotlib')
hv.Curve(target_curve_HOPG).opts(logy=True)* hv.Scatter(target_curve_HOPG.iloc[peaks]).opts(color='k', marker='s').opts(logy=True)

# +
peaks_1ML,_ =  sp.signal.find_peaks(target_curve_1ML,height = 1E-5, prominence =  1E-6)
pk_1ML_width = sp.signal.peak_widths(target_curve_1ML, peaks_1ML, rel_height=0.5)

hv.extension('bokeh')
#hv.extension('matplotlib')
hv.Curve(
    target_curve_HOPG).opts(
    logy = True,
    color = 'black')* hv.Scatter(
    target_curve_HOPG.iloc[peaks]).opts(
    color = 'black',
    fill_color = 'white',
    marker = 'triangle',
    size = 10).opts(
    logy = True)*hv.Curve(
    target_curve_1ML).opts(
    logy = True,color='blue')* hv.Scatter(
    target_curve_1ML.iloc[peaks_1ML]).opts(
    color = 'blue',
    fill_color = 'white',
    marker = 'triangle',
    size = 10).opts(
    logy=True)

# + [markdown] tags=[]
# ### HOPG  & 1ML peaks 
# #### find peaks unig sp.find_peaks
# -

grid_3D_gap_sg = savgolFilter_xr(grid_3D_gap, polyorder= 3, window_length= 11)
grid_3D_gap_sg['layer'] = layer_N.label
####
#grid_3D_gap_sg_1deriv =  grid_3D_gap_sg.differentiate(coord = 'bias_mV')
#grid_3D_gap_sg_2deriv = savgolFilter_xr(grid_3D_gap_sg_1deriv).differentiate(coord = 'bias_mV')
#grid_3D_gap_sg_2deriv_sg_dps = find_peaks_xr(-1*savgolFilter_xr(grid_3D_gap_sg_2deriv))
#grid_3D_gap_sg_2deriv_sg_dps_pad = peak_pad(grid_3D_gap_sg_2deriv_sg_dps)
#grid_3D_gap_sg_2deriv_sg_pks_pad

# +
# choose the data channel with groupby 

grid_3D_gap_STS = xr.merge([grid_3D_gap_sg[['I_fb','LDOS_fb','CBM_I_mV','VBM_I_mV','CBM_LIX_mV','VBM_LIX_mV','layer']] , grid_3D_gap_sg_2deriv_sg_dps_pad[['LDOS_fb_peaks']]]).dropna(dim = 'X')
#grid_3D_gap_STS.groupby('layer')['1ML'].mean(dim='stacked_X_Y').LDOS_fb.plot()
layer_STS_P = grid_3D_gap_STS.groupby('layer').mean(dim='stacked_X_Y')

# how select layers in xarray
#layer_STS_P[dict (layer = [0])]
STS_P_1ML = layer_STS_P.sel(layer = '1ML').LDOS_fb
STS_P_2ML = layer_STS_P.sel(layer = '2ML').LDOS_fb
STS_P_3ML = layer_STS_P.sel(layer = '3ML').LDOS_fb
STS_P_4ML = layer_STS_P.sel(layer = '4ML').LDOS_fb
STS_P_HOPG = layer_STS_P.sel(layer = 'HOPG').LDOS_fb

# +
print(grid_3D_gap_STS.groupby('layer')['1ML'].mean(dim='stacked_X_Y').CBM_LIX_mV.values,
      grid_3D_gap_STS.groupby('layer')['1ML'].mean(dim='stacked_X_Y').VBM_LIX_mV.values)
print(grid_3D_gap_STS.groupby('layer')['2ML'].mean(dim='stacked_X_Y').CBM_LIX_mV.values,
      grid_3D_gap_STS.groupby('layer')['2ML'].mean(dim='stacked_X_Y').VBM_LIX_mV.values)
print(grid_3D_gap_STS.groupby('layer')['3ML'].mean(dim='stacked_X_Y').CBM_LIX_mV.values,
      grid_3D_gap_STS.groupby('layer')['3ML'].mean(dim='stacked_X_Y').VBM_LIX_mV.values)
print(grid_3D_gap_STS.groupby('layer')['4ML'].mean(dim='stacked_X_Y').CBM_LIX_mV.values,
      grid_3D_gap_STS.groupby('layer')['4ML'].mean(dim='stacked_X_Y').VBM_LIX_mV.values)


print(grid_3D_gap_STS.groupby('layer')['1ML'].mean(dim='stacked_X_Y').CBM_I_mV.values,
      grid_3D_gap_STS.groupby('layer')['1ML'].mean(dim='stacked_X_Y').VBM_I_mV.values)
print(grid_3D_gap_STS.groupby('layer')['2ML'].mean(dim='stacked_X_Y').CBM_I_mV.values,
      grid_3D_gap_STS.groupby('layer')['2ML'].mean(dim='stacked_X_Y').VBM_I_mV.values)
print(grid_3D_gap_STS.groupby('layer')['3ML'].mean(dim='stacked_X_Y').CBM_I_mV.values,
      grid_3D_gap_STS.groupby('layer')['3ML'].mean(dim='stacked_X_Y').VBM_I_mV.values)
print(grid_3D_gap_STS.groupby('layer')['4ML'].mean(dim='stacked_X_Y').CBM_I_mV.values,
      grid_3D_gap_STS.groupby('layer')['4ML'].mean(dim='stacked_X_Y').VBM_I_mV.values)

# +
STS_P_1ML_1deriv = STS_P_1ML.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()
STS_P_1ML_2deriv = STS_P_1ML_1deriv.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()

STS_P_3ML_1deriv = STS_P_3ML.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()
STS_P_3ML_2deriv = STS_P_3ML_1deriv.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()


STS_P_4ML_1deriv = STS_P_4ML.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()
STS_P_4ML_2deriv = STS_P_4ML_1deriv.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()


hv.extension('bokeh')
pn.Row(hv.Curve(hv.Dataset(STS_P_1ML)).opts(axiswise = True, fig_size = 20).relabel('LDOS'),pn.Column(
          hv.Curve(hv.Dataset(STS_P_1ML_1deriv)).opts(axiswise = True, fig_size = 20).relabel('d(LDOS)/dV'),
          hv.Curve(hv.Dataset(STS_P_1ML_2deriv)).opts(axiswise = True, fig_size = 20).relabel('d2(LDOS)/dV2')))
#STS_P_1ML.plot()
#STS_P_1ML_1deriv.plot()
#.plot()
# -

# ##### HOPG,  STS peaks that is not in gap

grid_3D_gap_STS.groupby('layer')['1ML'].mean(dim='stacked_X_Y').CBM_LIX_mV.values

# +
STS_P_HOPG_1deriv = STS_P_HOPG.differentiate(coord = 'bias_mV').rolling(bias_mV = 10, center = True).mean()
STS_P_HOPG_2deriv = STS_P_HOPG_1deriv.differentiate(coord = 'bias_mV').rolling(bias_mV = 10, center = True).mean()


STS_P_HOPG_2deriv_dps,_ =  sp.signal.find_peaks(-1*STS_P_HOPG_2deriv,prominence=2E-18)
# use the prominence to choose the proper peak position 
STS_P_HOPG_2deriv_dps_width = sp.signal.peak_widths(STS_P_HOPG_2deriv, STS_P_HOPG_2deriv_dps, rel_height=0.5)

STS_P_HOPG_2deriv_dps_filtered = STS_P_HOPG_2deriv_dps[
    (STS_P_HOPG_2deriv.bias_mV[STS_P_HOPG_2deriv_dps].values
     > grid_3D_gap_STS.groupby('layer')['1ML'].mean(dim='stacked_X_Y').CBM_LIX_mV.values )
    +(STS_P_HOPG_2deriv.bias_mV[STS_P_HOPG_2deriv_dps].values <
      grid_3D_gap_STS.groupby('layer')['1ML'].mean(dim='stacked_X_Y').VBM_LIX_mV.values)]

# dip positions in 2nd derivative 
STS_P_HOPG_2deriv_dps_plt = hv.Curve(
    hv.Dataset(STS_P_HOPG_2deriv)).relabel('d2(LDOS)/dV2').opts(
    axiswise = True)*hv.Scatter(
    hv.Dataset(STS_P_HOPG_2deriv[STS_P_HOPG_2deriv_dps])).opts(
    color= 'k', marker="cross",size = 10).opts( 
    axiswise = True)
# dip positions apply to LDOS 
STS_P_HOPG_pks_plt = hv.Curve(
    hv.Dataset(STS_P_HOPG)).opts(axiswise = True).relabel('LDOS')*hv.Scatter(
    hv.Dataset(STS_P_HOPG[STS_P_HOPG_2deriv_dps_filtered])).opts(
    color= 'k', marker="cross",size = 10, axiswise = True).relabel('LDOS')

STS_P_HOPG_2deriv_dps_plt.relabel('d2(LDOS)/dV2').opts(axiswise = True) + STS_P_HOPG_pks_plt.opts(axiswise = True)

# + [markdown] tags=[]
# #### 1ML peak detection 

# +
STS_P_1ML_1deriv = STS_P_1ML.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()
STS_P_1ML_2deriv = STS_P_1ML_1deriv.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()


STS_P_1ML_2deriv_dps,_ =  sp.signal.find_peaks(-1*STS_P_1ML_2deriv,prominence=2E-18)
# use the prominence to choose the proper peak position 
STS_P_1ML_2deriv_dps_width = sp.signal.peak_widths(STS_P_1ML_2deriv, STS_P_1ML_2deriv_dps, rel_height=0.5)

STS_P_1ML_2deriv_dps_filtered = STS_P_1ML_2deriv_dps[
    (STS_P_1ML_2deriv.bias_mV[STS_P_1ML_2deriv_dps].values
     > grid_3D_gap_STS.groupby('layer')['1ML'].mean(dim='stacked_X_Y').CBM_LIX_mV.values )
    +(STS_P_1ML_2deriv.bias_mV[STS_P_1ML_2deriv_dps].values <
      grid_3D_gap_STS.groupby('layer')['1ML'].mean(dim='stacked_X_Y').VBM_LIX_mV.values)]

# dip positions in 2nd derivative 
STS_P_1ML_2deriv_dps_plt = hv.Curve(
    hv.Dataset(STS_P_1ML_2deriv)).relabel('d2(LDOS)/dV2').opts(
    axiswise = True)*hv.Scatter(
    hv.Dataset(STS_P_1ML_2deriv[STS_P_1ML_2deriv_dps])).opts(
    color= 'k', marker="cross",size = 10).opts( 
    axiswise = True)
# dip positions apply to LDOS 
STS_P_1ML_pks_plt = hv.Curve(
    hv.Dataset(STS_P_1ML)).opts(axiswise = True).relabel('LDOS')*hv.Scatter(
    hv.Dataset(STS_P_1ML[STS_P_1ML_2deriv_dps_filtered])).opts(
    color= 'k', marker="cross",size = 10, axiswise = True).relabel('LDOS')

STS_P_1ML_2deriv_dps_plt.relabel('d2(LDOS)/dV2').opts(axiswise = True) + STS_P_1ML_pks_plt.opts(axiswise = True)

# + [markdown] tags=[]
# #### 2ML peak detection 

# +
STS_P_2ML_1deriv = STS_P_2ML.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()
STS_P_2ML_2deriv = STS_P_2ML_1deriv.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()


STS_P_2ML_2deriv_dps,_ =  sp.signal.find_peaks(-1*STS_P_2ML_2deriv,prominence=2E-18)
# use the prominence to choose the proper peak position 
STS_P_2ML_2deriv_dps_width = sp.signal.peak_widths(STS_P_2ML_2deriv, STS_P_2ML_2deriv_dps, rel_height=0.5)

STS_P_2ML_2deriv_dps_filtered = STS_P_2ML_2deriv_dps[
    (STS_P_2ML_2deriv.bias_mV[STS_P_2ML_2deriv_dps].values
     > grid_3D_gap_STS.groupby('layer')['2ML'].mean(dim='stacked_X_Y').CBM_LIX_mV.values )
    +(STS_P_2ML_2deriv.bias_mV[STS_P_2ML_2deriv_dps].values <
      grid_3D_gap_STS.groupby('layer')['2ML'].mean(dim='stacked_X_Y').VBM_LIX_mV.values)]

# dip positions in 2nd derivative 
STS_P_2ML_2deriv_dps_plt = hv.Curve(
    hv.Dataset(STS_P_2ML_2deriv)).relabel('d2(LDOS)/dV2').opts(
    axiswise = True)*hv.Scatter(
    hv.Dataset(STS_P_2ML_2deriv[STS_P_2ML_2deriv_dps])).opts(
    color= 'k', marker="cross",size = 10).opts( 
    axiswise = True)
# dip positions apply to LDOS 
STS_P_2ML_pks_plt = hv.Curve(
    hv.Dataset(STS_P_2ML)).opts(axiswise = True).relabel('LDOS')*hv.Scatter(
    hv.Dataset(STS_P_2ML[STS_P_2ML_2deriv_dps_filtered])).opts(
    color= 'k', marker="cross",size = 10, axiswise = True).relabel('LDOS')

STS_P_2ML_2deriv_dps_plt.relabel('d2(LDOS)/dV2').opts(axiswise = True) + STS_P_2ML_pks_plt.opts(axiswise = True)

# + [markdown] jp-MarkdownHeadingCollapsed=true tags=[]
# #### 3ML peak detection 
# -

grid_3D_gap_STS.groupby('layer')['3ML'].mean(dim='stacked_X_Y').VBM_LIX_mV.values

# + jp-MarkdownHeadingCollapsed=true tags=[]
STS_P_3ML_1deriv = STS_P_3ML.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()
STS_P_3ML_2deriv = STS_P_3ML_1deriv.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()


STS_P_3ML_2deriv_dps,_ =  sp.signal.find_peaks(-1*STS_P_3ML_2deriv,prominence=2E-18)
# use the prominence to choose the proper peak position 
STS_P_3ML_2deriv_dps_width = sp.signal.peak_widths(STS_P_3ML_2deriv, STS_P_3ML_2deriv_dps, rel_height=0.5)

STS_P_3ML_2deriv_dps_filtered = STS_P_3ML_2deriv_dps[
    (STS_P_3ML_2deriv.bias_mV[STS_P_3ML_2deriv_dps].values
     > grid_3D_gap_STS.groupby('layer')['3ML'].mean(dim='stacked_X_Y').CBM_LIX_mV.values )
    +(STS_P_3ML_2deriv.bias_mV[STS_P_3ML_2deriv_dps].values <
      grid_3D_gap_STS.groupby('layer')['3ML'].mean(dim='stacked_X_Y').VBM_LIX_mV.values)]

# dip positions in 2nd derivative 
STS_P_3ML_2deriv_dps_plt = hv.Curve(
    hv.Dataset(STS_P_3ML_2deriv)).relabel('d2(LDOS)/dV2').opts(
    axiswise = True)*hv.Scatter(
    hv.Dataset(STS_P_3ML_2deriv[STS_P_3ML_2deriv_dps_filtered])).opts(
    color= 'k', marker="cross",size = 10).opts( 
    axiswise = True)
# dip positions apply to LDOS 
STS_P_3ML_pks_plt = hv.Curve(
    hv.Dataset(STS_P_3ML)).opts(axiswise = True).relabel('LDOS')*hv.Scatter(
    hv.Dataset(STS_P_3ML[STS_P_3ML_2deriv_dps_filtered])).opts(
    color= 'k', marker="cross",size = 10, axiswise = True).relabel('LDOS')

STS_P_3ML_2deriv_dps_plt.relabel('d2(LDOS)/dV2').opts(axiswise = True) + STS_P_3ML_pks_plt.opts(axiswise = True)
# -

# #### 4ML peak detection

# +
STS_P_4ML_1deriv = STS_P_4ML.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()
STS_P_4ML_2deriv = STS_P_4ML_1deriv.differentiate(coord = 'bias_mV').rolling(bias_mV = 3, center = True).mean()


STS_P_4ML_2deriv_dps,_ =  sp.signal.find_peaks(-1*STS_P_4ML_2deriv,prominence=2E-18)
# use the prominence to choose the proper peak position 
STS_P_4ML_2deriv_dps_width = sp.signal.peak_widths(STS_P_4ML_2deriv, STS_P_4ML_2deriv_dps, rel_height=0.5)

STS_P_4ML_2deriv_dps_filtered = STS_P_4ML_2deriv_dps[
    (STS_P_4ML_2deriv.bias_mV[STS_P_4ML_2deriv_dps].values
     > grid_3D_gap_STS.groupby('layer')['4ML'].mean(dim='stacked_X_Y').CBM_LIX_mV.values )
    +(STS_P_4ML_2deriv.bias_mV[STS_P_4ML_2deriv_dps].values <
      grid_3D_gap_STS.groupby('layer')['4ML'].mean(dim='stacked_X_Y').VBM_LIX_mV.values)]

# dip positions in 2nd derivative 
STS_P_4ML_2deriv_dps_plt = hv.Curve(
    hv.Dataset(STS_P_4ML_2deriv)).relabel('d2(LDOS)/dV2').opts(
    axiswise = True)*hv.Scatter(
    hv.Dataset(STS_P_4ML_2deriv[STS_P_4ML_2deriv_dps])).opts(
    color= 'k', marker="cross",size = 10).opts( 
    axiswise = True)
# dip positions apply to LDOS 
STS_P_4ML_pks_plt = hv.Curve(
    hv.Dataset(STS_P_4ML)).opts(axiswise = True).relabel('LDOS')*hv.Scatter(
    hv.Dataset(STS_P_4ML[STS_P_4ML_2deriv_dps_filtered])).opts(
    color= 'k', marker="cross",size = 10, axiswise = True).relabel('LDOS')


STS_P_4ML_2deriv_dps_plt.relabel('d2(LDOS)/dV2').opts(axiswise = True) + STS_P_4ML_pks_plt.opts(axiswise = True)

# +
STS_P_1ML_pks_plt = hv.Curve(
    hv.Dataset(STS_P_1ML), label = '1ML').opts(axiswise = True)*hv.Scatter(
    hv.Dataset(STS_P_1ML[STS_P_1ML_2deriv_dps_filtered])).opts(
     marker="cross",size = 10, axiswise = True)

STS_P_2ML_pks_plt = hv.Curve(
    hv.Dataset(STS_P_2ML), label = '2ML').opts(axiswise = True)*hv.Scatter(
    hv.Dataset(STS_P_2ML[STS_P_2ML_2deriv_dps_filtered])).opts(
     marker="cross",size = 10, axiswise = True)


STS_P_3ML_pks_plt = hv.Curve(
    hv.Dataset(STS_P_3ML), label = '3ML').opts(axiswise = True)*hv.Scatter(
    hv.Dataset(STS_P_3ML[STS_P_3ML_2deriv_dps_filtered])).opts(
     marker="cross",size = 10, axiswise = True)

STS_P_4ML_pks_plt = hv.Curve(
    hv.Dataset(STS_P_4ML), label = '4ML').opts(axiswise = True)*hv.Scatter(
    hv.Dataset(STS_P_4ML[STS_P_4ML_2deriv_dps_filtered])).opts(
     marker="cross",size = 10, axiswise = True)


(STS_P_1ML_pks_plt*STS_P_2ML_pks_plt*STS_P_3ML_pks_plt*STS_P_4ML_pks_plt)

# +
hv.Curve(hv.Dataset(STS_P_1ML)).opts(axiswise = True, fig_size = 20).relabel('LDOS')*hv.Scatter(hv.Dataset(STS_P_1ML[peaks_STS_P_1ML_2deriv])).opts(color= 'k',s =20, axiswise = True, fig_size = 20).relabel('LDOS')



#STS_P_1ML_2deriv[peaks_STS_P_1ML_2deriv]
# -

STS_P_1ML_2deriv.iloc[peaks_STS_P_1ML_2deriv]

STS_HOPG = grid_3D_gap_STS.groupby('layer')['HOPG'].mean(dim='stacked_X_Y').LDOS_fb
STS_1ML = grid_3D_gap_STS.groupby('layer')['1ML'].mean(dim='stacked_X_Y').LDOS_fb
STS_P_2ML = grid_3D_gap_STS.groupby('layer')['2ML'].mean(dim='stacked_X_Y').LDOS_fb
STS_P_3ML = grid_3D_gap_STS.groupby('layer')['3ML'].mean(dim='stacked_X_Y').LDOS_fb
STS_P_4ML = grid_3D_gap_STS.groupby('layer')['4ML'].mean(dim='stacked_X_Y').LDOS_fb

peaks_1ML,_ =  sp.signal.find_peaks(STS_1ML,height = 1E-5, prominence =  1E-6)
pk_1ML_width = sp.signal.peak_widths(STS_1ML, peaks_1ML, rel_height=0.5)

# +
grid_3D_gap_STS = xr.merge([grid_3D_gap_sg[['I_fb','LDOS_fb','layer']] , grid_3D_gap_sg_2deriv_sg_dps_pad[['LDOS_fb_peaks']]]).dropna(dim = 'X').to_dataframe()
# LDOS data + 2deriv dps
#grid_3D_gap_STS[]
#grid_3D_gap_STS
#[grid_3D_gap_STS.layer == 'HOPG']

target_curve_HOPG =grid_3D_gap_STS[grid_3D_gap_STS.layer == 'HOPG'].LDOS_fb.groupby('bias_mV').mean()
target_curve_1ML =grid_3D_gap_STS[grid_3D_gap_STS.layer == '1ML'].LDOS_fb.groupby('bias_mV').mean()

# -

savgolFilter_xr(target_curve_HOPG)

# +
hv.extension('bokeh')
#hv.extension('matplotlib')
hv.Curve(
    target_curve_HOPG).opts(
    logy = True,
    color = 'black')* hv.Scatter(
    target_curve_HOPG.iloc[peaks]).opts(
    color = 'black',
    fill_color = 'white',
    marker = 'triangle',
    size = 10)


#grid_3D_gap_STS.groupby(['layer']).mean()
#xr.concat ([grid_3D_gap_sg[['I_fb','LDOS_fb','layer']], grid_3D_gap_sg_2deriv_sg_pks_pad[['LDOS_fb_peaks']]], dim=["X"])
# -

#grid_3D_gap_sg_2deriv_sg_pks_pad.LDOS_fb_peaks
grid_3D_gap_sg_2deriv_sg_pks_pad.LDOS_fb.mean('X').plot()


# + tags=[]
peaks_1ML,_ =  sp.signal.find_peaks(target_curve_1ML,height = 1E-5, prominence =  1E-6)
pk_1ML_width = sp.signal.peak_widths(target_curve_1ML, peaks_1ML, rel_height=0.5)

#hv.extension('bokeh')
hv.extension('matplotlib')
hv.Curve(
    target_curve_HOPG).opts(
    logy = True,
    color = 'black')* hv.Scatter(
    target_curve_HOPG.iloc[peaks]).opts(
    color = 'black',
    marker = '|', 
    s = 200).opts(
    logy = True)*hv.Curve(
    target_curve_1ML).opts(
    logy = True,color='blue')* hv.Scatter(
    target_curve_1ML.iloc[peaks_1ML]).opts(
    color = 'blue',
    marker = '|',
    s = 200).opts(
    logy=True,
    fontsize= 'x-large', 
    fig_size= 200)
# -

# ### HOPG  & 1ML peaks 
# #### find peaks using 2nd derivatives 
#
#

target_curve_1ML = lineSTS_P[lineSTS_P.layer == '1ML'].LDOS_fb.groupby('bias_mV').mean()

# + [markdown] tags=[] jp-MarkdownHeadingCollapsed=true
# ### Pyramid  2, 3, 4 ML STS 
# #### find peaks unig sp.find_peaks

# +
target_curve_2ML = lineSTS_P[lineSTS_P.layer == '2ML'].LDOS_fb.groupby('bias_mV').mean()
target_curve_3ML = lineSTS_P[lineSTS_P.layer == '3ML'].LDOS_fb.groupby('bias_mV').mean()
target_curve_4ML = lineSTS_P[lineSTS_P.layer == '4ML'].LDOS_fb.groupby('bias_mV').mean()

target_curve_2ML = target_curve_2ML* 1E9
target_curve_2ML.index  = target_curve_2ML.index*1E-3
target_curve_2ML.index.names  = ['Bias (V)']

target_curve_3ML = target_curve_3ML* 1E9
target_curve_3ML.index  = target_curve_3ML.index*1E-3
target_curve_3ML.index.names  = ['Bias (V)']

target_curve_4ML = target_curve_4ML* 1E9
target_curve_4ML.index  = target_curve_4ML.index*1E-3
target_curve_4ML.index.names  = ['Bias (V)']

peaks_2ML,_ =  sp.signal.find_peaks(target_curve_2ML,height = 1E-5, prominence =  1E-6)
pk_2ML_width = sp.signal.peak_widths(target_curve_2ML, peaks_2ML, rel_height=0.5)
peaks_3ML,_ =  sp.signal.find_peaks(target_curve_3ML,height = 1E-5)#, prominence =  1E-7)
pk_3ML_width = sp.signal.peak_widths(target_curve_3ML, peaks_3ML, rel_height=0.5)
peaks_4ML,_ =  sp.signal.find_peaks(target_curve_4ML,height = 1E-5)#, prominence =  1E-7)
pk_4ML_width = sp.signal.peak_widths(target_curve_4ML, peaks_4ML, rel_height=0.5)


#hv.extension('bokeh')
hv.extension('matplotlib')
hv.Curve(
    target_curve_2ML).opts(
    color='brown') *hv.Scatter(
    target_curve_2ML.iloc[peaks_2ML]).opts(
    color='brown', marker="|", s = 200)*hv.Curve(
    target_curve_3ML).opts(
    color='orange')*hv.Scatter(
    target_curve_3ML.iloc[peaks_3ML]).opts(
    color='orange', marker="|", s = 200)*hv.Curve(
    target_curve_4ML).opts(
    color='gold')*hv.Scatter(
    target_curve_4ML.iloc[peaks_4ML]).opts(
    color='gold', marker="|", s = 200).opts(
    logy= True,
    ylabel  = 'dI/dV (nA/V)',
    xlabel = 'Bias (V)',
    fontsize= 'x-large', 
    fig_size= 200)



#hv.Curve(target_curve_1ML).opts(logy=True)* hv.Scatter(target_curve_1ML.iloc[peaks_1ML]).opts(color='k', marker='s').opts(logy=True)

# -

hv.extension('matplotlib')
hv.Curve(target_curve_HOPG).opts(
    color = 'black')* hv.Curve(
    target_curve_1ML).opts(
    color = 'blue').opts(
    logy= True,
    ylabel  = 'dI/dV (nA/V)',
    xlabel = 'Bias (V)',
    fontsize= 'x-large', 
    fig_size= 200)

layer_N = z_profile_up_terrace.to_xarray() 
# pd.qcut 활용한 결과를 xarray로 바꿈. 

# +
#grid_3D_sg_2deriv_sg_pks_pad
grid_3D_gap_sg_1deriv = grid_3D_gap_sg.differentiate(coord = 'bias_mV')
grid_3D_gap_sg_1deriv_sg = savgolFilter_xr(grid_3D_gap_sg_1deriv, window_length= 11, polyorder= 3)
#grid_3D_gap_sg_1deriv.LDOS_fb.mean(['X','Y']).plot()
#grid_3D_gap_sg_1deriv_sg.LDOS_fb.isel(X = 10,  Y = 0).plot()

grid_3D_gap_sg_2deriv = grid_3D_gap_sg_1deriv_sg.differentiate(coord = 'bias_mV')
grid_3D_gap_sg_2deriv.LDOS_fb.isel(X = 10,  Y = 0).plot()

# d(LDOS)dV pks & dps
grid_3D_gap_sg_2deriv_dps = find_peaks_xr(-1*grid_3D_gap_sg_2deriv)
grid_3D_gap_sg_2deriv_dps_pad = peak_pad(grid_3D_gap_sg_2deriv_dps)
# -

grid_3D_gap_sg_2deriv_dps.LDOS_fb_peaks.isel(X = 1).values[0]

# +
grid_3D_gap_sg_2deriv_dps_pad_1 = grid_3D_gap_sg_2deriv_dps.LDOS_fb_peaks.isel(X = 1).values[0]
#grid_3D_gap_sg_2deriv_dps_pad.LDOS_fb_peaks_pad[grid_3D_gap_sg_2deriv_dps_pad_1]
grid_3D_gap_sg_2deriv_dps.LDOS_fb.isel(X = 1,bias_mV = grid_3D_gap_sg_2deriv_dps_pad_1)

#grid_3D_gap_sg_2deriv_dps_pad
# -

grid_3D_gap_sg_2deriv_dps.LDOS_fb.isel(X = 1).isel(bias_mV = grid_3D_gap_sg_2deriv_dps_pad_1)

grid_3D_gap_sg_2deriv_dps_pad_1

sns.lineplot(x = target_curve.bias_mV,y =target_curve  )
sns.scatte

rplot(x = target_curve.bias_mV[peaks],y =target_curve[peaks]  )

# +

sns.set_style("whitegrid")
z_profile_plot = sns.lineplot(x = "X", y = "topography", data = z_profile)
z_profile_plot.set_xlabel ("X (nm)")
z_profile_plot.set_ylabel ("Z (nm)")
#z_profile.topography.plot()
# -


grid_3D['layers'] = Nlayers
grid_3D

Nlayers.drop ("index")
Nlayers.set_coords ("X")
Nlayers.set_dims ("X")
#Nlayers.set_coords (names = )
#Nlayers.set_coords(grid_xr.X.values)
#Nlayers.coords = grid_xr.X


# +
#sliderX_v = grid_3D.X[sliderX.value].values
#sliderY_v = grid_3D.Y[sliderY.value].values

##


peak_detection_curve  = grid_3D_sg_2deriv_sg_pks_pad.isel(X = sliderX.value, Y = sliderY.value)
#peak_detection_curve.LIX_fb[peak_detection_curve.LIX_fb_peaks_pad.values.astype(int)]
#peak_detection_curve.LIX_fb_peaks_pad.values.dropna().astype(int)
peak_detection_curve.LIX_fb_peaks.values
# -

# ### object --> array 변환 필요 

grid_3D_sg_2deriv_sg_pks_pad

# ### peak --> histogram --> peak existing area --> ROI selection 
# ####  특정한 peak 이 존재하는 위치 ROI
# ###3 or 특정한 peak 의 intensity FWHM plot  ==> topography  image correlation ==> physical meaning --> defect --> peak 을 enhance 

# +
#grid_3D_sg_2deriv_sg_pks_pad.LIX_fb_peaks_pad.to_dataframe()
grid_3D_sg_2deriv_sg_pks.LIX_fb_peaks.to_dataframe().explode('LIX_fb_peaks')

grid_3D_sg_2deriv_sg_pks_pad.LIX_fb_peaks_pad.to_dataframe().unstack()

# -

# # plot curves & peak points 
#

# + [markdown] tags=[]
# # Data link between curves 

# +
peak_detection_curve.variables.mapping.keys()
data_ch = 'LIX_fb'

target_curve = peak_detection_curve[data_ch]
#target_curve.bias_mV

peaks,_ =  sp.signal.find_peaks(target_curve)
pk_width = sp.signal.peak_widths(target_curve,peaks, rel_height=0.5)

hv.Curve(target_curve)*hv.Scatter(target_curve[peaks]).opts(color='k', marker='s')



# +
sns.lineplot(x = target_curve.bias_mV,y =target_curve  )
sns.scatterplot(x = target_curve.bias_mV[peaks],y =target_curve[peaks]  )

#target_curve[peaks]

#target_curve.bias_mV

# +
sns.lineplot(x = target_curve.bias_mV,y =target_curve  )


#sns.lineplot(data = peak_detection_curve, x ='bias_mV', y = 'LIX_fb')

#peak_detection_curve
#pk = peak_detection_curve.LIX_fb_peaks

# + [markdown] tags=[]
# # Histogram for peaks 

# + [markdown] tags=[]
# ## 3. Live Interactive plot (?)

# + active=""
#
# -

I_fb_hvplot = grid_3D.I_fb.hvplot(groupby='bias_mV', cmap = 'bwr', aspect = 'equal', width = 400, widget_type='scrubber', widget_location='bottom')
# usnig hvplot with widget 
LIX_fb_hvplot = grid_3D.LIX_fb.hvplot(groupby='bias_mV', cmap = 'bwr', aspect = 'equal', width = 400, widget_type='scrubber', widget_location='bottom')
#I_fb_hvplot+LIX_fb_hvplot => plot link 하는 방법을 찾아야함. 
(I_fb_hvplot+LIX_fb_hvplot)

# + [markdown] jp-MarkdownHeadingCollapsed=true tags=[]
# # Manual Piecewised Affine Transform

# +
topo2Dxr = sxmData_PF_Sat
topo3Dxr = gridData_topo_PF_Sat

from pathlib import Path
import tkinter as tk
from tkinter import messagebox
from tkinter import filedialog
import tkinter.font as tkFont
import cv2
from PIL import Image
from PIL import ImageTk    
from csv import *

def btncmd_SltRemove():

    try:
        index_start = listbox_points.curselection()[0]
        index_end = listbox_points.curselection()[-1]
        listbox_points.delete(index_start,index_end)
    except:
        messagebox.showwarning("Alert", "Choose points in Listbox above.")
        return

    # print('index',index_start,index_end)
    for i in range(index_start,index_end+1):
        ptsList1.pop(index_start)
        ptsList2.pop(index_start)

    __reloadImage__()
    # print(ptsList1,ptsList2)
def btncmd_AllRemove():
    global ptsList1,ptsList2
    listbox_points.delete(0,listbox_points.size()-1)
    # print('index',index_start,index_end)
    ptsList1 = []
    ptsList2 = []
    
    __reloadImage__()

def event_slt1(event):
    # global img_2Dtopo, ptsList1, ptsList2
    if len(ptsList1) > len(ptsList2):
        messagebox.showwarning("Alert", "Choose point in right side.")
        return


    ptsList1.append([event.x,event.y])

    __reloadImage__()

    # print(event.x,event.y)


    if len(ptsList1) == len(ptsList2):
        listbox_points.insert(tk.END, str([ptsList1[-1],ptsList2[-1]]))
def event_slt2(event):
    # global img_3Dtopo,ptsList1, ptsList2

    if len(ptsList1) < len(ptsList2):
        messagebox.showwarning("Alert", "Choose point in left side.")
        return

    ptsList2.append([event.x,event.y])

    __reloadImage__()


    if len(ptsList1) == len(ptsList2):
        listbox_points.insert(tk.END, str([ptsList1[-1],ptsList2[-1]]))

        
def btncmd_loadPts():
    with open(str(Path().absolute())+"\pts_for_PAF.csv","r") as file:
        Reader = reader(file,delimiter=',')
        
        loadList1 = []
        loadList2 = []
        
        
        
        line = 0
        for row in Reader:
            if row == []:
                continue
            if line == 0:# 2D pts
            
                for i in range(len(row)):
                    ptsNum = len(row)
                    newPts1 = [int(row[i][1:-1].split(',')[0]), int(row[i][1:-1].split(',')[1])]
                    if newPts1 in ptsList1:
                        messagebox.showwarning("Alert", "Same point(s) already loaded")
                        return
                    else:
                        ptsList1.append( newPts1)
                        loadList1.append(newPts1)
                line +=1
                
            elif line == 1:# 3D pts
                
                for i in range(len(row)):
                    newPts2 = [int(row[i][1:-1].split(',')[0]), int(row[i][1:-1].split(',')[1])]
                    ptsList2.append( newPts2)
                    loadList2.append(newPts2)
                break
                
    for i in range(len(loadList1)):
        listbox_points.insert(tk.END, str([loadList1[i],loadList2[i]]))
    
    __reloadImage__()
    messagebox.showinfo("Information", str(ptsNum)+"Points Loaded.")
    # fpath = filedialog.askopenfilename(title="abcdefu",filetypes = (".csv","*.csv"))
    # filePathStr.set(fpath)

def btncmd_savePts():
    with open(str(Path().absolute())+"\pts_for_PAF.csv","w") as file:
        Writer = writer(file,delimiter=',')
        Writer.writerows([ptsList1,ptsList2])
        messagebox.showinfo("Information", "Points Saved.")
    
    

def __reloadImage__():
    img_2Dcircle = np.zeros((np.shape(img_2Dtopo)[0],np.shape(img_2Dtopo)[0],3),np.uint8)
    for i in range(len(ptsList1)):
        color = (127,127,127)
        img_2Dcircle = cv2.circle(img_2Dcircle,(ptsList1[i][0],ptsList1[i][1]),5,color,2)

    img_2Dres = img_2Dtopo + img_2Dcircle
    imgArr_2Dres = Image.fromarray(img_2Dres)
    imgtk_2Dres = ImageTk.PhotoImage(image=imgArr_2Dres)

    label_2Dtopo.configure(image = imgtk_2Dres)
    label_2Dtopo.image = imgtk_2Dres


    img_3Dcircle = np.zeros((np.shape(img_3Dtopo)[0],np.shape(img_3Dtopo)[0],3),np.uint8)
    for i in range(len(ptsList2)):
        color = (127,127,127)
        img_3Dcircle = cv2.circle(img_3Dcircle,(ptsList2[i][0],ptsList2[i][1]),5,color,2)

    img_3Dres = img_3Dtopo + img_3Dcircle
    imgArr_3Dres = Image.fromarray(img_3Dres)
    imgtk_3Dres = ImageTk.PhotoImage(image=imgArr_3Dres)

    label_3Dtopo.configure(image = imgtk_3Dres)
    label_3Dtopo.image = imgtk_3Dres
    
    
    
"""
Data load
"""      

# Dimension, Spacing 뽑아내기
Xdim2D, Ydim2D = topo2Dxr.coords.dims['X'],topo2Dxr.coords.dims['Y']
Xdim3D, Ydim3D = topo3Dxr.coords.dims['X'],topo3Dxr.coords.dims['Y']
Xspacing2D, Xspacing3D = topo2Dxr.X_spacing, topo3Dxr.X_spacing

# Topography들만 뽑아내기
topo2D = topo2Dxr.z_fwd_df.values
topo3D = topo3Dxr.topography.values

# 3D의 경우 2D의 scale에 맞도록 scaling

sizeRatio = int(Xspacing3D / Xspacing2D)

# print(Xspacing2D,Xspacing3D,sizeRatio)
topo3D_resized = resize(topo3D, 
                       ( sizeRatio*Xdim3D , sizeRatio* Ydim3D), 
                       anti_aliasing=True)

# ref_topo3D_resized = cv2.resize( ref_topo3D, 
#                                ( sizeRatio*Xdim3D , sizeRatio* Ydim3D), 
#                                interpolation=cv2.INTER_AREA)

# 3D의 경우 2D의 크기에 맞게 padding
topo3D_pad = np.empty((Xdim2D, Ydim2D))
topo3D_pad[int(Xdim2D/2) - int(len(topo3D_resized)/2) : int(Xdim2D/2) + int(len(topo3D_resized)/2), \
                  int(Ydim2D/2) - int(len(topo3D_resized[0])/2) : int(Ydim2D/2) + int(len(topo3D_resized[0])/2)]\
            = topo3D_resized

# Feature finding을 위한 grayscaling
topo2D_gray = grayscaling(topo2D)
topo3D_gray = grayscaling(topo3D_pad)

"""
UI part
"""

root = tk.Tk()
root.resizable(False, False)




ptsList1 = []
ptsList2 = []
# src = cv2.imread(r"D:\CloudStation\Data set production (scripts)\testIMG.png")


font = tkFont.Font(family="Arial", size=40, weight="bold", slant="italic")
label_title = tk.Label(text="   Feature Match Manual   ",font=font,relief="ridge",)
label_title.pack(side="top")

img_2Dtopo = cv2.applyColorMap(topo2D_gray.astype('uint8'), cv2.COLORMAP_INFERNO)


imgArr_2Dtopo = Image.fromarray(img_2Dtopo)
imgtk_2Dtopo = ImageTk.PhotoImage(image=imgArr_2Dtopo)

label_2Dtopo = tk.Label(image = imgtk_2Dtopo)
label_2Dtopo.image = imgtk_2Dtopo
label_2Dtopo.pack(side = "left")




img_3Dtopo = cv2.applyColorMap(topo3D_gray.astype('uint8'), cv2.COLORMAP_INFERNO)


imgArr_3Dtopo = Image.fromarray(img_3Dtopo)
imgtk_3Dtopo = ImageTk.PhotoImage(image=imgArr_3Dtopo)

label_3Dtopo = tk.Label(image = imgtk_3Dtopo)
label_3Dtopo.image = imgtk_3Dtopo
label_3Dtopo.pack(side = "right")


frameList = tk.Frame(root)
frameList.pack()

scrollbar = tk.Scrollbar(frameList)
scrollbar.pack(side="right",fill="y")

listbox_points = tk.Listbox(frameList,selectmode="extended",height = 18,yscrollcommand = scrollbar.set)
listbox_points.pack(side="top")



frameRemove = tk.Frame(root)
frameRemove.pack()
btn_removePts = tk.Button(frameRemove,text="Remove Pair(s)",command = btncmd_SltRemove)
btn_removePts.pack(side="left",pady=10,ipady=10)
btn_removeAll = tk.Button(frameRemove,text="Remove All",command = btncmd_AllRemove)
btn_removeAll.pack(side="right",pady=10,ipady=10)


filePathStr = tk.StringVar()
frameSaveLoad = tk.Frame(root)
frameSaveLoad.pack()
btn_savePts = tk.Button(frameSaveLoad,text="Save Pair(s)",command = btncmd_savePts)
btn_savePts.pack(side="left",pady=5,ipady=5)
btn_loadPts = tk.Button(frameSaveLoad,text="Load Pair(s)",command = btncmd_loadPts)
btn_loadPts.pack(side="right",pady=5,ipady=5)



btn_quit = tk.Button(text="Quit",command = root.destroy)
btn_quit.pack(side="bottom",pady=10,ipadx=10,ipady=10)




label_2Dtopo.bind("<Button-1>",event_slt1)
label_3Dtopo.bind("<Button-1>",event_slt2)


root.mainloop()
print("Points 1:",ptsList1)
print("Points 2:",ptsList2)

pts1_PAF = ptsList1
pts2_PAF = ptsList2


# + [markdown] tags=[]
# # 3. Piecewise-Affine-Transform
# -

def piecewiseAffineByPts(topo2Dxr,topo3Dxr,pts2D,pts3D):
# Import modules
    from skimage.transform import PiecewiseAffineTransform, warp
    

    
    # Dimension, Spacing 뽑아내기
    Xdim2D, Ydim2D = topo2Dxr.coords.dims['X'],topo2Dxr.coords.dims['Y']
    Xdim3D, Ydim3D = topo3Dxr.coords.dims['X'],topo3Dxr.coords.dims['Y']
    Xspacing2D, Xspacing3D = topo2Dxr.X_spacing, topo3Dxr.X_spacing

    # Topography들만 뽑아내기
    topo2D = topo2Dxr.z_fwd_df.values
    topo3D = topo3Dxr.topography.values

    # 3D의 경우 2D의 scale에 맞도록 scaling

    sizeRatio = int(Xspacing3D / Xspacing2D)

    # print(Xspacing2D,Xspacing3D,sizeRatio)
    topo3D_resized = resize(topo3D, 
                           ( sizeRatio*Xdim3D , sizeRatio* Ydim3D), 
                           anti_aliasing=True)

    # ref_topo3D_resized = cv2.resize( ref_topo3D, 
    #                                ( sizeRatio*Xdim3D , sizeRatio* Ydim3D), 
    #                                interpolation=cv2.INTER_AREA)

    # 3D의 경우 2D의 크기에 맞게 padding
    topo3D_pad = np.empty((Xdim2D, Ydim2D))
    topo3D_pad[int(Xdim2D/2) - int(len(topo3D_resized)/2) : int(Xdim2D/2) + int(len(topo3D_resized)/2), \
                      int(Ydim2D/2) - int(len(topo3D_resized[0])/2) : int(Ydim2D/2) + int(len(topo3D_resized[0])/2)]\
                = topo3D_resized
    
    # Definition of images sizes
    rows, cols = topo2D.shape[0], topo2D.shape[1]

    # Source and destination coordinates
    src = np.array(pts2D)
    dst = np.array(pts3D)

    # 조각 아핀 선언하기, source-destination 순서를 반대로 입력해야 잘 나온다.
    tform = PiecewiseAffineTransform()
    tform.estimate(src, dst)

    # 출력 크기 정하고, 아핀 실행하기
    topo3D_PAF = warp(topo3D_pad, tform, output_shape=(rows, cols))

    # 그래프로 확인
    fig, ax = plt.subplots(1, 3, figsize = (15, 5))
    
    isns.imshow(topo2D, ax = ax[0], cmap = cmap_topo, robust = True, perc = (0, 99.5))
    ax[0].set_title('2D topo')
    
    isns.imshow(topo3D_pad, ax = ax[1], cmap = cmap_topo, robust = True, perc = (0, 99.5))
    ax[1].plot(dst[:, 0],dst[:, 1], '.b')
    ax[1].set_title('3D topo - Before PAF')


    isns.imshow(topo3D_PAF, ax = ax[2], cmap = cmap_topo, robust = True, perc = (0, 99.5))
    ax[2].set_title('3D topo - After PAF')
    ax[2].plot(src[:, 0],src[:, 1], '.b')

    plt.show()
    
    
    
    
    result_xr = xr.Dataset(
        {
            "topography" : (["X","Y"], topo3D_PAF),
        },
        coords = {
            "X": (["X"], topo2Dxr.X.values), 
            "Y": (["Y"], topo2Dxr.Y.values), 
        }
    )
    result_xr.attrs["title"] = topo2Dxr.title
    
    
    result_xr.attrs['image_size']= topo2Dxr.image_size
    result_xr.attrs['X_spacing']= topo2Dxr.X_spacing
    result_xr.attrs['Y_spacing']= topo2Dxr.Y_spacing  
    result_xr.attrs['freq_X_spacing']= 1/topo2Dxr.X_spacing
    result_xr.attrs['freq_Y_spacing']= 1/topo2Dxr.Y_spacing
    
    return result_xr,tform
# gridData_topo_PF_Sat_AF_PAF, tform = piecewiseAffineByPts(sxmData_PF_Sat, gridData_topo_PF_Sat_AF,pts1_PAF,pts2_PAF)
gridData_topo_PF_Sat_AF_PAF, tform = piecewiseAffineByPts(sxmData_PF_Sat, gridData_topo_PF_Sat,pts1_PAF,pts2_PAF) # by manual
isns.imshow(gridData_topo_PF.topography.values,cmap = cmap_topo)

# + [markdown] tags=[]
# # FFT Analysis 
#

# +
# TwoD_fft_xr 사용 


# -

# # Streamz package 사용 
#

# +



#slider = ipw.IntSlider(description='time', min=0, max=10)
grid_3D_LIX_fb_test = grid_3D.LIX_fb.interactive().sel(bias_mV=pnw.DiscreteSlider).hvplot(cmap = 'bwr', aspect = 'equal',width=400)
stream_xy = hv.streams.Tap(source = grid_topo_bokeh, x = np.nan, y = np.nan)
grid_3D_LIX_fb_test


# -

stream_xy


# + jp-MarkdownHeadingCollapsed=true tags=[]
#grid_topo.coords
grid_topo_bokeh = grid_topo.hvplot( aspect='equal').opts(tools=['box_select', 'lasso_select'])
#topography_map = hv.Image(grid_topo.topography.values)

# Create the Tap stream with the points element as the source
# We set the x and y here with starting values
stream = hv.streams.Tap(source = grid_topo_bokeh, x = np.nan, y = np.nan)



def cross_lines (stream_tapXY):
    point_stream = grid_topo.sel(X = [stream.x], Y =  [stream.y], method = "nearest")
    #x = point_stream.X
    x = point_stream.X.values[0]
    y = point_stream.Y.values[0]
    return hv.VLine(x) * hv.HLine(y)

#cross_line_sel = hv.DynamicMap(cross_lines, streams=stream) 

grid_topo_bokeh*cross_lines(stream)

# +
#grid_topo_bokeh = grid_topo.hvplot( aspect='equal').opts(tools=['box_select', 'lasso_select'])
#topography_map = hv.Image(grid_topo.topography.values)

# Create the Tap stream with the points element as the source
# We set the x and y here with starting values
stream = hv.streams.Tap(source = grid_topo_bokeh, x = np.nan, y = np.nan)

point_stream = grid_topo.sel(X = [stream.x], Y =  [stream.y], method = "nearest")
#x = point_stream.X
x = point_stream.X.values[0]
y = point_stream.Y.values[0]
grid_topo_bokeh*hv.VLine(x)* hv.HLine(y)

# +


data_channel = list(grid_3D.data_vars.keys())
# data channel names =>  to select 

slider_mV = pnw.IntSlider(name='bias_mV', start=0, end=len(grid_3D.bias_mV))
slider_x = pnw.IntSlider(name='X', start = 0, end = len(grid_3D.X))
slider_y = pnw.IntSlider(name='Y', start = 0, end = len(grid_3D.Y))


dmap = hv.DynamicMap(grid_3D[data_channel].isel(X = slider_x , Y = slider_y, bias_mV= slider_mV ), kdims=['X','Y','bias_mV'])
dmap = dmap.redim.values(Data_channel=data_channel)
dmap

# -

len(grid_3D.bias_mV)
#list(grid_3D.data_vars.keys())

# + jp-MarkdownHeadingCollapsed=true tags=[]
# make a function that displays the location when called.
def location(x, y):
    """Display pane showing the x and y values"""
    return pn.pane.Str('Click at %0.3f, %0.3f' % (x, y), width = 200)
layout = pn.Row(grid_topo_bokeh, pn.bind(location, x = stream.param.x, y = stream.param.y))
# display the container
layout
# -

point_stream.X.values[0]

point_stream = grid_topo.sel(X = [stream.x], Y =  [stream.y], method = "nearest")
point_stream.X

 hv.Curve()
#grid_3D_gap.I_fb.
hv.Curve(grid_3D_gap.isel(X = x, Y = y), kdims='bias_mV')

# +
stream
def lissajous_crosshair(t, a=3, b=5, delta=np.pi/2):
    (x,y) = lissajous(t,a,b,delta)
    return hv.VLine(x) * hv.HLine(y)

crosshair = hv.DynamicMap(lissajous_crosshair, kdims='t').redim.range(t=(-3.,3.))

path = hv.Path(lissajous(lin))

# +

hv.Curve(grid_3D_gap.LDOS_fb.sel(X = [stream.x], Y =  [stream.y], method = "nearest"), kdims = "bias_mV")

#grid_3D_gap.LDOS_fb.sel( X = stream.x, Y = stream.y)

#hv.Curve(grid_3D_gap.LDOS_fb.sel(x = stream.x, y = stream.y))

# +
#Decare a holoview figure 
hv.extension('bokeh')
# Declare map to select 
topography_map = hv.Image(grid_topo.topography.values)

# Declare Tap stream with heatmap as source and initial values
posxy = hv.streams.Tap(source = topography_map, x='X', y='Y')



# Define function to compute STS  based on tap location
def tap_LDOS(x, y):
    return hv.Curve(grid_3D_gap.isel(X = x, Y = y), kdims='bias_mV',
                   label=f'X: {x}, Y: {y}')


# Connect the Tap stream to the tap_histogram callback
tap_dmap = hv.DynamicMap(tap_LDOS, streams=[posxy])
# Display the Heatmap and Curve side by side

topography_map+tap_dmap

# +
topography_map.array()

sliderX = pnw.IntSlider(name='X', 
                       start = 0 ,
                       end = grid_3D.X.shape[0]) 
sliderY = pnw.IntSlider(name='X', 
                       start = 0 ,
                       end = grid_3D.Y.shape[0]) 

XY = Stream.event('XY',x=sliderX,y=sliderY)
def marker(x,y):
    return hv.VLine(x) * hv.HLine(y)

image = topography_map
dmap = hv.DynamicMap(marker, streams=[XY()])

image * dmap
# -

#Decare a holoview figure 
hv.extension('bokeh')
# use the 'bokeh' extension for hv 
# hv.extension('matplotlib')
'''LDOS_fb_hv = grid_3D_gap.LDOS_fb.hvplot(groupby='bias_mV',
                                        aspect='equal',
                                        robust = True)

'''
# Declare map to select 
topography_map = hv.Image(grid_topo.topography.values)
topography_map.array()

# +
# Declare points as source of selection stream
selection = hv.streams.Selection1D(source = topography_map)

# Write function that uses the selection indices to slice points and compute stats
def selected_info(index):
    arr = topography_map.array()[index]
    if index:
        label = 'Mean x, y: %.3f, %.3f' % tuple(arr.mean(axis=0))
    else:
        label = 'No selection'
    return topography_map.clone(arr, label=label)#.opts(color='red')

# Combine points and DynamicMap
selected_points = hv.DynamicMap(selected_info, kdims = ['X','Y'])# streams=[selection])

layout = topography_map.opts(tools=['box_select', 'lasso_select']) #+ selected_points

layout
# -

grid_topo.topography.values[[0,1],[1,2]]

# +
hv.extension('bokeh')
# use the 'bokeh' extension for hv 
# hv.extension('matplotlib')
grid_3D_gap.LDOS_fb.hvplot(groupby='bias_mV',
                        aspect='equal',
                        robust = True).redim.range(LDOS_fb = tuple(np.array([grid_3D_gap.LDOS_fb.min().values,
                                                                          grid_3D_gap.LDOS_fb.max().values])))
# how to set  aspect ratop
# how to set cmap range ('tuple')



# Declare map to select 
topography_map = hv.Image(grid_topo.topography.values)

# Declare points as source of selection stream
selection = hv.streams.Selection1D(source = topography_map)



# Write function that uses the selection indices to slice points and compute stats
def selected_info(index):
    arr = topography_map.array()[index]
    if index:
        label = 'Mean x, y: %.3f, %.3f' % tuple(arr.mean(axis=0))
    else:
        label = 'No selection'
    return topography_map.clone(arr, label=label).opts(color='red')

# Combine points and DynamicMap
selected_points = hv.DynamicMap(selected_info, streams=[selection])
layout = topography_map.opts(tools=['box_select', 'lasso_select']) + selected_points

layout



# -


grid_3D_gap.LDOS_fb.hvplot(groupby='bias_mV',
                        aspect='equal',
                        robust = True)+grid_3D_gap.gap_size_I.hvplot(aspect='equal',
                                                                     robust = True)


# +
tuple(np.array([grid_3D_gap.dIdV.min().values, grid_3D_gap.dIdV.max().values]))


#grid_3D_gap.dIdV.max()


# +
#   plot two image at the same time 
grid_3D_gap.dIdV.hvplot(groupby='bias_mV',
                        width=500,
                        height=400,
                        robust = True)+grid_3D_gap.LDOS_fb.hvplot(groupby='bias_mV',
                                                                  width=500,
                                                                  height=400, 
                                                                  robust = True)




slider = pnw.IntSlider(name='bias_mV', 
                       start = 0 ,
                       end = grid_3D.bias_mV.shape[0]) 
# -
# ## t-test confidence interval calc. 


# +
#bound_box.bounds
x_bounds_msk = (grid_3D.X > bound_box.bounds[0] ) & (grid_3D.X < bound_box.bounds[2])
y_bounds_msk = (grid_3D.Y > bound_box.bounds[1] ) & (grid_3D.Y < bound_box.bounds[3])

grid_3D_bbox = grid_3D.where (grid_3D.X[x_bounds_msk] + grid_3D.Y[y_bounds_msk])

fig,axs = plt.subplots (nrows = 1,
                        ncols = 3,
                        figsize = (9,3))

isns.imshow(grid_3D.I_fb.sel(bias_mV = 0 ),
            ax =  axs[0],
            robust = True)
# add rectangle for bbox 

isns.imshow(grid_3D_bbox.I_fb.sel(bias_mV = 0 ),
            ax =  axs[1],
            robust = True)
sns.lineplot(x = "bias_mV",
             y = "LIX_fb", 
             data = grid_3D_bbox.to_dataframe(),
             ax = axs[2])

plt.show()



# plot STS at the selected points 
# use the seaborn (confident interval : 95%) 
# sns is figure-level function 


# +
#slct_pts
pts = grid_channel_hv_points.iloc[slct_pts.index].dframe().set_index(['X', 'Y'])

pts_xr = xr.Dataset.from_dataframe(pts)

grid_3D_slct_pts = xr.combine_by_coords ([grid_3D, pts_xr], compat = 'override', join = 'inner')
#y_pts = points.iloc[slct_pts.index].dframe().Y
#grid_3D.sel(X = x_pts,Y = y_pts)
#grid_3D.I_fb.isel(bias_mV = 0).plot()

fig, axs = plt.subplots(ncols = 2, nrows = 1, figsize = (8,4))

grid_3D_slct_pts.I_fb.plot(ax = axs[0], robust = True) 

sns.lineplot(x = "bias_mV",            
             y = "LIX_fb", 
             data = grid_3D_slct_pts.to_dataframe(),
             ax = axs[1])
plt.show
#grid_3D_slct_pts
#
#sns.relplot(x="bias_mV",
#            y="LIX_fb", 
#            kind="line",
#            data=grid_3D_slct_pts.to_dataframe())
# check. sn.relplot is  figure-level function

# +
#grid_3D_gap.hvplot(groupby='bias_mV') + grid_3D_gap.hvplot.line(x='Y', y= 'I_fb')

sliderX = pnw.IntSlider(name='X', 
                       start = 0 ,
                       end = grid_3D.X.shape[0]) 
sliderY = pnw.IntSlider(name='X', 
                       start = 0 ,
                       end = grid_3D.Y.shape[0]) 

const_HLine =  grid_3D_gap.Y[60].values

from holoviews import opts
grid_3D_gap.I_fb.hvplot(groupby='bias_mV',
                        x = 'X', 
                        y = 'Y',
                        width=500,
                        height=400,
                        robust = True) *hv.HLine(const_HLine).opts(opts.HLine(color='blue', line_width=3))

#grid_3D_gap.I_fb.hvplot(groupby='bias_mV',x = 'X', y = 'Y', width=500, height=400, robust = True) + grid_3D_gap.I_fb.hvplot.line(x='bias_mV' ,width=400, height=400)

# +

slider = pnw.IntSlider(name='bias_mV', 
                       start = 0 ,
                       end = grid_3D.bias_mV.shape[0]) 
# use the end value as a bias_mV idx max value
grid_3D.LIX_fb.interactive(width=800).isel(bias_mV = slider).plot(cmap = 'bwr',robust=True)




# +
hv.extension('bokeh')

#hv.extension('matplotlib')
grid_3D.I_fb.interactive(width=400).sel(bias_mV=pnw.DiscreteSlider).plot( cmap='bwr',robust=True)


# + [markdown] tags=[]
# # Smoohting & rolling 
# -

def find_peaks_xr(xrdata): 
    from scipy.signal import find_peaks
    xrdata_prcssd = xrdata.copy(deep = True)
    print('Find peaks in STS to an xarray Dataset.')
    xAxis = xrdata.X.size
    yAxis = xrdata.Y.size
    for data_ch in xrdata:
        if len(xrdata[data_ch].dims)==2:
            # smoothing filter only for the 3D data set
            pass
            
        else :
            print (data_ch)
            """xrdata_prcssd[data_ch+'_peaks']= xr.DataArray(np.ones((xAxis,yAxis), dtype = object),
                                                             dims=["X", "Y"],
                                                             coords={"X": xrdata.X, "Y": xrdata.Y} )"""
            xrdata_prcssd[data_ch+'_peaks'] = xr.DataArray (
                np.array([ find_peaks(xrdata[data_ch].isel(X = x, Y = y).values)[0] 
                          for x in range(xAxis)  
                          for y in range(yAxis)], dtype = object ).reshape(xAxis,yAxis),
                dims=["X", "Y"],
                coords={"X": xrdata.X, "Y": xrdata.Y})
            
                    # use the "xrdata_prcssd[data_ch].values[x,y,:]"  for "xrdata_prcssd" 
                    # not ".isel(X = x, Y = y).values"""

    return xrdata_prcssd
grid_3D_sg_pks = find_peaks_xr(grid_3D_sg)



# +
grid_3D_gap.LDOS_fb[20,20,:].plot()
grid_3D_gap_roll = grid_3D_gap.rolling(bias_mV =5,center=True).mean()
grid_3D_gap_roll.LDOS_fb[20,20,:].plot()
grid_3D_gap_sg.LDOS_fb[20,20,:].plot()

# Need a interactive plot for comparison 






# +
# choose the data_variables 

data_vars_pf_list = ['I_fb', 'LIX_fb', 'LIX_fb_offst', 'LIX_fb_offst_CB',  'LIX_fb_offst_VB']
grid_3D_fb_gap = grid_3D_fb_gap[data_vars_pf_list]
grid_3D_fb_gap = grid_3D_fb_gap.rolling(bias_mV = 3, center = True).mean()

# channel 골라내기 

pf_deg = 13
pf_coeff = grid_3D_fb_gap.polyfit('bias_mV', deg = pf_deg)

pf_coeff_names = list(pf_coeff.keys())
#print(pf_coeff_names)
for i, data_ch in enumerate(grid_3D_fb_gap): 
    print(data_ch)
    grid_3D_fb_gap[data_ch].bias_mV.dropna(data_ch)
bias_mV_pf
# -

grid_3D_fb_gap.LIX_fb_offst_VB.isel(X=40, Y=40)

    for pf_coeff_name in pf_coeff_names:
        if pf_coeff_name.startswith(data_ch):
            print(pf_coeff_name)
            # bias_mV_range  지정해줘야함. 
            grid_3D_fb_gap[data_ch+'_pf'] =  xr.polyval(bias_mV_pf,
                                                pf_coeff[pf_coeff_name],
                                                degree_dim='degree')
        else : # if  data vars are not consistent
            pass #print(data_ch , 'is not start with', pf_coeff_name[i] )


#grid_3D_fb_gap.LIX_fb_offst_CB_pf.isel(X=40,Y=40).plot()
#grid_3D_fb_gap.LIX_fb_offst_CB.isel(X=40,Y=40).plot()
grid_3D_fb_gap.LIX_fb_offst_VB_pf.isel(X=40,Y=40).plot()
#grid_3D_fb_gap.LIX_fb_offst_VB.isel(X=40,Y=40).plot()

# +

grid_3D_fb_gap_pf


LIX_CB_xr_roll  = LIX_CB_xr.rolling(bias_mV = 3, center = True).mean()
LIX_CB_xr_roll

pf_deg = 13
pf_coeff = LIX_CB_xr_roll.polyfit('bias_mV', deg = pf_deg)

pf_coeff_names = list(pf_coeff.keys())
print(pf_coeff_names)

# +

for pf_coeff_name in pf_coeff_names:
    if pf_coeff_name.startswith(LIX_CB):
        print(pf_coeff_name)
        LIX_CB_xr_roll['pf'] =  xr.polyval(LIX_CB_xr_roll.bias_mV,
                                   pf_coeff[pf_coeff_name],
                                   degree_dim='degree') 
    else : # if  data vars are not consistent
        print(data_ch , 'is not start with', pf_coeff_name[i] )

# +
pf_coeff_name = list(pf_coeff.keys())
print(pf_coeff_name)
for i, data_ch in enumerate(grid_3D_fb): 
    print(data_ch)
    if pf_coeff_name[i].startswith(data_ch):
        print(pf_coeff_name[i])
        grid_3D_fb_roll[data_ch+'_pf'] =  xr.polyval(grid_3D_fb.bias_mV,
                                                pf_coeff[pf_coeff_name[i]],
                                                degree_dim='degree') 
    else : # if  data vars are not consistent
        print(data_ch , 'is not start with', pf_coeff_name[i] )




# how to find partial data set fo polyfit. 


# smoothing by rolling 
grid_3D_fb_gap_roll = grid_3D_fb_gap.rolling(bias_mV = 3, center = True).mean().dropna('bias_mV')
pf_deg = 13
pf_coeff = grid_3D_fb_gap_roll.polyfit('bias_mV', deg = pf_deg)
# -

pf_coeff_name = list(pf_coeff.keys())
print(pf_coeff_name)
for i, data_ch in enumerate(grid_3D_fb): 
    print(data_ch)
    if pf_coeff_name[i].startswith(data_ch):
        print(pf_coeff_name[i])
        grid_3D_fb_roll[data_ch+'_pf'] =  xr.polyval(grid_3D_fb.bias_mV,
                                                pf_coeff[pf_coeff_name[i]],
                                                degree_dim='degree') 
    else : # if  data vars are not consistent
        print(data_ch , 'is not start with', pf_coeff_name[i] )

# +
# smoothing by rolling 
grid_3D_fb_roll = grid_3D_fb_prcssd.rolling(bias_mV = 3, center = True).mean().dropna('bias_mV')
pf_deg = 13
pf_coeff = grid_3D_fb_roll.polyfit('bias_mV',
                              deg = pf_deg)
pf_coeff_name = list(pf_coeff.keys())
print(pf_coeff_name)
for i, data_ch in enumerate(grid_3D_fb): 
    print(data_ch)
    if pf_coeff_name[i].startswith(data_ch):
        print(pf_coeff_name[i])
        grid_3D_fb_roll[data_ch+'_pf'] =  xr.polyval(grid_3D_fb.bias_mV,
                                                pf_coeff[pf_coeff_name[i]],
                                                degree_dim='degree') 
    else : # if  data vars are not consistent
        print(data_ch , 'is not start with', pf_coeff_name[i] )

grid_3D_fb = grid_3D_fb.rolling(bias_mV = 3, center = True).mean().dropna('bias_mV')
# rolling 3 mean & drop_na 


# for polyfits 
# polyfit_degree = 

pf_deg = 13
pf_coeff = grid_3D_fb.polyfit('bias_mV',
                              deg = pf_deg)
pf_coeff_name = list(pf_coeff.keys())
print(pf_coeff_name)
for i, data_ch in enumerate(grid_3D_fb): 
    print(data_ch)
    if pf_coeff_name[i].startswith(data_ch):
        print(pf_coeff_name[i])
        grid_3D_fb[data_ch+'_pf'] =  xr.polyval(grid_3D_fb.bias_mV,
                                                pf_coeff[pf_coeff_name[i]],
                                                degree_dim='degree') 
    else : # if  data vars are not consistent
        print(data_ch , 'is not start with', pf_coeff_name[i] )
# -



grid_3D_gap_sg

# # interactive plot with Holoviews
#

# +
# interactive plot
import holoviews as hv
from holoviews import opts
# %matplotlib qt5

hv.extension('matplotlib')

opts.defaults(opts.Scatter3D(color='Value', cmap='fire', edgecolor='black', s=50))


# +
hv_ds = hv.Dataset(grid_3D_fb.LIX_fb)
print(hv_ds)

XY_bias = hv_ds.to(hv.Image, kdims=["Y", "X"], dynamic= True)
XY_bias.opts(colorbar=True, fig_size=100)
# X,Y dimension direction  Trnaspose 필요함. 


# + [markdown] tags=[]
# *    when you need to call the datavarible names  in xr data 
# >  keys =  list (grid_3D_fb.keys())[0]

# +
grid_3D_fb = grid_3D_fb.rolling(bias_mV = 3, center = True).mean().dropna('bias_mV')
# rolling 3 mean & drop_na 


# for polyfits 
# polyfit_degree = 

pf_deg = 13
pf_coeff = grid_3D_fb.polyfit('bias_mV',
                              deg = pf_deg)
pf_coeff_name = list(pf_coeff.keys())
print(pf_coeff_name)
for i, data_ch in enumerate(grid_3D_fb): 
    print(data_ch)
    if pf_coeff_name[i].startswith(data_ch):
        print(pf_coeff_name[i])
        grid_3D_fb[data_ch+'_pf'] =  xr.polyval(grid_3D_fb.bias_mV,
                                                pf_coeff[pf_coeff_name[i]],
                                                degree_dim='degree') 
    else : # if  data vars are not consistent
        print(data_ch , 'is not start with', pf_coeff_name[i] )
# -


grid_3D_fb
grid_3D_fb.isel(X=1,Y=1).plot.scatter(x = 'bias_mV', y = 'LIX_fb')
grid_3D_fb.isel(X=1,Y=1).plot.scatter(x = 'bias_mV', y = 'LIX_fb_pf')


# +
test_keys = pd.Series(grid_3D_fb.keys()) 
pf_chs = pd.Series(grid_3D_fb.polyfit('bias_mV',deg = 7).keys())
pf_chs[0].startswith(test_keys[0])


for data_ch in grid_3D_fb: 
    print(data_ch)
    
    
        data_ch+pf_coeff_name[0]

pf_chs = list(grid_3D_fb.polyfit('bias_mV',deg = 7).keys())
# polyfit coefficient + deg value 


grid_3D_fb['I_fb_pf'] = 

fit = xr.polyval(grid_3D_fb.bias_mV, grid_3D_fb.polyfit('bias_mV',deg = 7))
grid_3D_fb.polyfit('bias_mV',deg = 7).I_fb_polyfit_coefficients.isel( X=1, Y=1)
# position 별 coefficient 
grid_3D_fb.

test_keys = grid_3D_fb.keys()
list(test_keys)
# pyty fit

#grid_3D

#grid_3D.isel(X=1,Y=1).LIX_fwd.plot()
#grid_3D.isel(X=1,Y=1).LIX_bwd.plot()
#grid_3D.isel(X=1,Y=1).I_fb.plot()


# +
#grid_3D_fb.polyfit('bias_mV',deg = 7).keys()



test_keys[0] 

list(grid_3D_fb.polyfit('bias_mV',deg = 7).keys())

fit = xr.polyues(grid_3D_fb.bias_mV, grid_3D_fb.polyfit('bias_mV',deg = 7))
# -

grid_3D_fb.I_fb.polyfit('bias_mV',deg = 23)
I_fb_pf_test  = xr.polyval(grid_3D_fb.bias_mV, grid_3D_fb.I_fb.polyfit('bias_mV',deg = 23), degree_dim='degree')
#I_fb_pf_test.polyfit_coefficients
I_fb_pf_test.polyfit_coefficients.isel(X=1,Y=1).plot()
grid_3D_fb.I_fb.isel(X=1,Y=1).plot()
Plot_xr_r(PlaneFitY_xr(grid_topo))



# grid_topo=xr.dataset(grid_topo)


# + [markdown] tags=[] jp-MarkdownHeadingCollapsed=true
# # ROI based on threshold 
#
# * using Otsu -> groupby 
# * using ROI to BBOX with opening & closing 

# +
grid_3D_th = hv.Dataset(ThrshldMultiOtsu_xr(grid_3D).LIX_unit_calc)


dmap = grid_3D_th.to(hv.Image,
                     kdims = dmap_plane,
                     dynamic = True )
dmap.opts(colorbar = True,
          cmap = 'bwr',
          frame_width = 400,
          aspect = 'equal').relabel('XY plane slicing: ')

# -

grid_3D_gap




# +
grid_3D_bbox = grid_3D.where (grid_3D.X[x_bounds_msk] + grid_3D.Y[y_bounds_msk])



test_3d = hv.Dataset(grid_3D.I_fb)
test_3d_0 = hv.Dataset(grid_3D.I_fb.isel(bias_mV =0))

test_dmap = test_3d.to(hv.Image, kdims=["X","Y"], dynamic = True )
test_dmap.opts(colorbar = True, cmap = 'bwr', frame_width=200,  aspect = 'equal') .relabel('XY plane slicing: ')


slct_pts = hv.streams.Selection1D(source=points)
test_dmap*points

# +
#slct_pts
pts = points.iloc[slct_pts.index].dframe().set_index(['X', 'Y'])

pts_xr = xr.Dataset.from_dataframe(pts)

grid_3D_slct_pts = xr.combine_by_coords ([grid_3D, pts_xr], compat = 'override', join = 'inner')
#y_pts = points.iloc[slct_pts.index].dframe().Y
#grid_3D.sel(X = x_pts,Y = y_pts)

sns.relplot(x="bias_mV", y="LIX_fb", kind="line",data=grid_3D_slct_pts.to_dataframe());


# + jp-MarkdownHeadingCollapsed=true tags=[]




############
# for image overview (realspace)
#Plot_xr_r(spmdata_xr, figsize= (5,6))

# #%matplotlib inline
#Plot_xr_r(ThrshldMultiOtsu_xr(spmdata_xr, multiclasses= 3))
#spmdata_xr.z_fwd_df.where(ThrshldMultiOtsu_xr(spmdata_xr, multiclasses= 3).z_fwd_df ==2).plot()

#Plot_xr_r(PlaneFitX_xr(PlaneFitY_xr(spmdata_xr)))

##########
# for histogram
#isns.imghist(PlaneFitX_xr(PlaneFitY_xr(spmdata_xr)).z_fwd_df.values, cmap = 'bone')

##########
# for FFt result plot 
#xr_fft_isns_plot(twoD_fft_xr( plane_x_fit_xr(plane_y_fit_xr(spmdata_xr))))
#spmdata_xr
#file_list_df.file_nameb
#spmdata_xr
# -

# # Plane Fit with Threshold ROI selection 
# * select ROI area with MultiOtsu result 
#     * Choose the area (0,1,2) 
#     * PlaneFitY with given area 
#     * replace the processed (Plane Fit) channels
# * Plot  "spmdata_xr_processed" 
#

# +
spmdata_xr.z_fwd_df.where(ThrshldMultiOtsu_xr(spmdata_xr, multiclasses= 3).z_fwd_df ==2).plot()
#  MultiOtsu selection (use "where") in  z_fwd_df 
# remove mean (axis=1) along Y direction 
#(spmdata_xr.z_fwd_df - spmdata_xr.z_fwd_df.where(ThrshldMultiOtsu_xr(spmdata_xr).z_fwd_df ==2).mean(axis =1)).plot()
isns.imshow(spmdata_xr.z_fwd_df - spmdata_xr.z_fwd_df.where(ThrshldMultiOtsu_xr(spmdata_xr).z_fwd_df ==2).mean(axis =1), robust = True)

spmdata_xr_processed = spmdata_xr.copy()
spmdata_xr_processed.z_fwd_df.values = spmdata_xr.z_fwd_df - spmdata_xr.z_fwd_df.where(ThrshldMultiOtsu_xr(spmdata_xr).z_fwd_df ==2).mean(axis =1)
spmdata_xr_processed.z_bwd_df.values = spmdata_xr.z_bwd_df - spmdata_xr.z_bwd_df.where(ThrshldMultiOtsu_xr(spmdata_xr).z_bwd_df ==2).mean(axis =1)


#########
# plot & save fig after processing 
spmdata_xr_processed_fig = Plot_xr_r(spmdata_xr_processed)
spmdata_xr_processed_fig.savefig(
    spmdata_xr_processed.title.split('\n')[0] +'_r.png')

spmdata_xr_processed_fft_fig = Plot_xr_fft(TwoD_FFT_xr(spmdata_xr))
spmdata_xr_processed_fft_fig.savefig(
    spmdata_xr_processed.title.split('\n')[0] +'_fft.png')


# + [markdown] tags=[]
# #  3.  Summary pptx file  generation
#
# * 1. presentation file saved as 'sxm_file_name_analysis.pptx'
# * 2. Title page 
# * 3. analsysis pages
#
#

# +
###############################
# new pptx file 
# overwrite everytime 
###############################


from pptx.util import Inches
from pptx.util import Cm, Pt

folder_name  = os.getcwd().split('\\')[-1]
# final folder 
pptx_file_name =  spmdata_xr.title.split('\n')[0]+ '_analysis.pptx'

prs_sxm = Presentation() # empty_file #빈파일 하나 지정 
#``` Page 0 ``````````
sld0_slide_layout = prs_sxm.slide_layouts[2] # Section Header slide ?(layout_type =  [2] )

slide0 = prs_sxm.slides.add_slide(sld0_slide_layout)
# add title slide in the pptx file # 슬라이드를 하나 더했음. 
sld0_title = slide0.shapes.title
sld0_title.text =spmdata_xr.title.split('\n')[0] + '\n' +  'in '+ folder_name 

sld0_subtitle = slide0.placeholders[1]
sld0_subtitle.text = 'Summary'



#``` Page 1 ``````````

sld1_slide_layout = prs_sxm.slide_layouts[5] # title + blank spaces
slide1 = prs_sxm.slides.add_slide(sld1_slide_layout)
slide1_title  = slide1.shapes.title
slide1_title.text = 'Topography and LDOS' # set page title # 제목설정 


left = top = Inches(1)
width = Inches(6)
height = Inches(1)

slide1_textbox = slide1.shapes.add_textbox(left, 
                                           top, 
                                           width, 
                                           height) # add text box 

paragraph0 = slide1_textbox.text_frame.add_paragraph()
# To control the font, use the paragraph
paragraph0.text = 'After Threshold ROI  Plane Fit' # 내용설정
paragraph0.level = 1   
paragraph0.font.size = Pt(14)
# adjust the font size
paragraph1 = slide1_textbox.text_frame.add_paragraph()
paragraph1.text = spmdata_xr.title
paragraph1.level = 2   
paragraph1.font.size = Pt(12)
# add paragraph in the contents  : file name & scan setting

# save figures  
# from Plot_xr_r & Plot_xr_fft

# check the image size
spmdata_xr_processed_fig.savefig(
    spmdata_xr_processed.title.split('\n')[0] +'_r.png',
    figsize= (5,6),
    bbox_inches='tight',
    dpi = 300)
spmdata_xr_processed_fft_fig.savefig(
    spmdata_xr_processed.title.split('\n')[0] +'_fft.png',
    figsize= (5,6),
    bbox_inches='tight',
    dpi = 300)
########
# use the add_picture to add above figures 
# check the imaage positions in ppt 
left = Inches(0)
top = Inches(3)

r_space_img = slide1.shapes.add_picture(
    spmdata_xr_processed.title.split('\n')[0] +'_r.png',
    left, 
    top)

left_2 = Inches(5)

k_space_img = slide1.shapes.add_picture(
    spmdata_xr_processed.title.split('\n')[0] +'_fft.png',
    left_2, 
    top)


prs_sxm.save(pptx_file_name)  

# add group_name & number of files as a paragraph
# -


spmdata_xr = Img2xr(file_list_df.file_name[0])
spmdata_xr

Plot_xr_r(grid_xr.isel(bias_mV = [0]))
#Plot_xr_r(spmdata_xr)

# + [markdown] tags=[]
# # Line Profile 
# * **Change the Backend Plot**
#     * eg) <u/> %matplotlib qt5<u>
#     * LineProfile_xr_GUI : for single channel 
#     * LineProfile2_xr_GUI : for two channels  
# * after analysis, save the output figure for pptxfile. 
#     * use the xrdata.title attributes as a file name 
#     

# +
# for example, get a line profile of z_fwd_df

# %matplotlib qt5

lp_start, lp_end, lp_fig = LineProfile_xr_GUI(
    spmdata_xr_processed, ch_N= 0 )
# saving figure of line profile 
lp_fig.savefig(spmdata_xr_processed.title.split('\n')[0]+'_lp_fig.png',figsize = (4,5))

# %matplotlib inline

# use the given (cx,ry) data points np array w/o GUI
'''
lp_start = np.array([2,1])
lp_end = np.array([100,100])
LineProfile_xr(spmdata_xr_processed, lp_start, lp_end)
'''


# + jp-MarkdownHeadingCollapsed=true tags=[]
# %matplotlib qt5

lp_start, lp_end, lp_fig1 ,lp_fig2 = LineProfile2_xr_GUI(
    spmdata_xr_processed, ch_N=[0,2])
#LineProfile2_xr_GUI(spmdata_xr_processed, ch_N=[0,2])

lp_fig1.savefig(spmdata_xr_processed.title.split('\n')[0]+'_lp_fig1.png')
lp_fig2.savefig(spmdata_xr_processed.title.split('\n')[0]+'_lp_fig2.png')



# %matplotlib inline

'''
# one can use the given (cx,ry) data points (np array w/o GUI)

lp_start = np.array([100,200])
lp_end = np.array([300,400])
lp_fig1 ,lp_fig2 = LineProfile2_xr(spmdata_xr_processed, lp_start, lp_end)
'''



# +
#####
# select I_fb & LIX_fb_offst for polyfit
#grid_3D_fb_gap
#
data_varis_list = list(grid_3D_fb_gap.data_vars.variables.mapping.keys())
data_varis_list.remove ('I_fb')
data_varis_list.remove('LIX_fb_offst')
data_varis_list                                                    

grid_3D_fb_pf = grid_3D_fb_gap.drop(data_varis_list)

# + [markdown] tags=[]
# # Add analysis page 
# * use the **AddAnalysisSlide** function 
#

# +
AddAnalysisSlidePage(pptx_file_name, 
                     'z_fwd line profile',
                    spmdata_xr_processed.title.split('\n')[0]+'_lp_fig.png')


AddAnalysisSlidePage(pptx_file_name, 
                     'z_fwd & LIX_fwd line profile',
                    spmdata_xr_processed.title.split('\n')[0]+'_lp_fig1.png')


AddAnalysisSlidePage(pptx_file_name, 
                     'z_fwd & LIX_fwd line profiles',
                    spmdata_xr_processed.title.split('\n')[0]+'_lp_fig2.png')
# -


# # more analysis will be added sooner or later 
# ## Espectially, the atomically resolved images + FFT analysis
# ## + Affine transform + Piecewise Affine transform, etc. 
#

#
# > 1. Use all the groups sxm files in folders? [default : Y]
# >> **OR**, 1-1. Choose the groups in 'sxm_file_groups' 
#     * file_list_df[file_list_df['group'] == sxm_file_groups[0]]
#
# > 2. include FFT? [default : Y]
# >> **OR**, 2-1. Choose ( include_fft == False)
#

# # xr data file save & open 

# +
# file save test
spmdata_xr.to_netcdf("HOPG_Wtip#12_LN2T_0036.nc")


# open saved file test
opentest = xr.open_dataset("HOPG_Wtip#12_LN2T_0036.nc")


# + [markdown] tags=[]
# ## pptx pages for each image 
# ###  page structure 
#
# | real space (**r**) |  real space (**r**)|  | momentum space (**k**) | momentum space (**k**)|
# | ---| ---|    | ---| ---|
# |topo (fwd)|topo(bwd)|    | topo_fft (fwd)|topo_fft (bwd)| 
# |LIX (fwd) |LIX (fwd)|    |  LIX_fft  (fwd)|   LIX_fft  bwd)|

# +
#bound_box.bounds
x_bounds_msk = (grid_3D.X > bound_box.bounds[0] ) & (grid_3D.X < bound_box.bounds[2])
y_bounds_msk = (grid_3D.Y > bound_box.bounds[1] ) & (grid_3D.Y < bound_box.bounds[3])

bbox = np.array(bound_box.bounds)
bbox[0:2]
fig, ax = plt.subplots()
isns.imshow(grid_3D.I_fb.sel(bias_mV =0 ),
            ax =  ax,
            robust = True)

rect = plt.Rectangle((bound_box.bounds[0], bound_box.bounds[1]) # left, bottom
                     , (bound_box.bounds[2]-bound_box.bounds[0]),  (bound_box.bounds[3]-bound_box.bounds[1]),
                     facecolor="white",
                     alpha=0.4)


ax.add_patch(rect)
plt.show()
# -

# # streams & interactive plot test ==> later 

# +
grid_3D_ref_img = hv.Image(grid_3D.I_fb.isel(bias_mV = 0))
# hv.Image 지정 
grid_channel_hv_coords =  hv.Points(grid_3D_ref_img)
stream = hv.streams.Tap(source=grid_channel_hv_coords, x=np.nan, y=np.nan)
# Connect the Tap stream to the tap_histogram callback
tap_dmap = hv.DynamicMap(grid_channel_hv_coords, streams=[stream])


# Display the Heatmap and Curve side by side
grid_3D_ref_img + tap_dmap
# -

.opts(
    opts.Curve(framewise=True, height=500, line_color='black', width=375, yaxis='right'),
    opts.Image( cmap='RdBu_r', fontsize={'xticks': '6pt'}, height=500,tools=['hover'], width=700, xrotation=90)
)

# +
grid_3D_ref_img = hv.Image(grid_3D.I_fb.isel(bias_mV = 0))
# hv.Image 지정 
grid_channel_hv_coords =  hv.Points(grid_3D_ref_img).opts(frame_width = 400,  
                                   aspect = 'equal', alpha = 0.2,                                   
                                   tools=['box_select', 'lasso_select','hover']
                                  )
stream = hv.streams.Tap(source=grid_channel_hv_coords, x=np.nan, y=np.nan)
# Connect the Tap stream to the tap_histogram callback
tap_dmap = hv.DynamicMap(grid_channel_hv_coords, streams=[stream])

grid_3D_ref_img*grid_channel_hv_coords+tap_dmap 


# +


# Connect the Tap stream to the tap_histogram callback
tap_dmap = hv.DynamicMap(grid_channel_hv_coords, streams=[stream])
tap_dmap 

#grid_channel_hv.closest([posxy])

#grid_channel_hv.select()

# +
grid_channel_hv_coords  = hv.Dataset(grid_3D_ref_img).relabel('for Point selection : ')
#  hv.Points 로 coord  position 지정 



img_coords = hv.Points(img, extents=extents)
labeled_img = img * img_coords * hv.Points([img.closest([(4.1,4.3)])]).opts(color='r')
img + labeled_img + img.sample([(4.1,4.3)])


hover_points = hv.Points(grid_channel_hv_image).opts(frame_width = 400,
                                                    color = 'k',
                                                    aspect = 'equal',
                                                    alpha = 0.1,                                   
                                                    tools=['hover'])

posxy  = hv.streams.Tap(source = hover_points, 
                        x = np.nan,
                        y = np.nan)

dmap*hover_points
# -

#posxy.x
#hvCurve(grid_channel_hv.select( X = posxy.x , Y =  posxy.x))
grid_channel_hv.select(( X = posxy.x ,  Y = posxy.x))



# +
def tap_STS(x, y):
    return hv.Curve(grid_channel_hv.select( X = x , Y = y),
                    kdims = 'bias_mV')#,label=f'X: {x}, Y: {y}')


tap_dmap = hv.DynamicMap(tap_STS, 
                    streams = [posxy])
(dmap*hover_points + tap_dmap)     

#
# -

posxy


# +
# Define function for STS based on tap location
def tap_STS(x, y):
    return hv.Curve(grid_channel_hv.select( X = x , Y = y),
                    kdims='bias_mV')#,label=f'X: {x}, Y: {y}')

tap_dmap = hv.DynamicMap(tap_STS, streams=[posxy])

#tap_dmap
# -

grid_channel_hv.map(X=0,Y=0)


# +
posxy_x = grid_3D.X.sel(X = [posxy.x],  method="nearest")
posxy_y = grid_3D.Y.sel(Y = [posxy.y],  method="nearest")

# Define function for STS based on tap location
def tap_STS(x, y):
    return hv.Curve(grid_channel_hv.select(X = posxy_x, Y = posxy_y), kdims='bias_mV')#,label=f'X: {x}, Y: {y}')

tap_dmap = hv.DynamicMap(tap_STS, streams=[posxy])

dmap+tap_dmap


# +
# Define function for STS based on tap location
def tap_STS(x, y):
    return hv.Curve(grid_channel_hv.select(X = posxy.x, Y = posxy.y), kdims='bias_mV')#,label=f'X: {x}, Y: {y}')


tap_dmap = hv.DynamicMap(tap_STS, streams=[posxy])

dmap+tap_dmap

# +
#grid_topo.coords
grid_topo_bokeh = grid_topo.hvplot( aspect='equal').opts(tools=['box_select', 'lasso_select'])
# tools option 

#topography_map = hv.Image(grid_topo.topography.values)

# Create the Tap stream with the points element as the source
# We set the x and y here with starting values
stream = hv.streams.Tap(source = grid_topo_bokeh, x = np.nan, y = np.nan)


# +
def peak_pad(xrdata):
    xrdata_prcssd = xrdata.copy(deep = True)
    xAxis = xrdata.X.size
    yAxis = xrdata.Y.size
    for ch in xrdata_prcssd:
        if ch.endswith('_peaks'):
            peaks = xrdata[ch].values
            peaks_count_max = max([ len(peaks_r_c) 
                            for peaks_r in peaks
                            for  peaks_r_c in peaks_r])
            padding_value = np.nan
            
            peaks_pad = np.array([
                np.pad(peaks_r_c.astype(float), 
                       (0,peaks_count_max-len(peaks_r_c)),
                       mode = 'constant', 
                       constant_values = padding_value)
                for peaks_r in peaks 
                for  peaks_r_c in peaks_r]).reshape((xAxis,yAxis,-1))
            
            xrdata_prcssd[ch+'_pad'] = xr.DataArray(peaks_pad, dims=["X", "Y","peaks"],
                coords={"X": xrdata.X, "Y": xrdata.Y, "peaks": np.arange(peaks_count_max)})
        else: pass
    return xrdata_prcssd
        
grid_3D_sg_pks_pad = peak_pad(grid_3D_sg_pks)
grid_3D_sg_pks_pad

# +
hv.extension('bokeh')
grid_3D_sg_LDOS_fft = hv.Dataset(np.log(TwoD_FFT_xr(grid_3D_sg).LIX_unit_calc_fft))
dmap_fft  = grid_3D_sg_LDOS_fft.to(hv.Image,
                                   kdims = ['freq_bias_mV', 'freq_Y'],#'freq_bias_mV', ['freq_X', 'freq_Y']
                                   dynamic = True)
#dmap_fft

dmap_fft.opts(colorbar = True,
              cmap = 'bwr',
              frame_width = 200,
              aspect = 'square').relabel('XY plane slicing FFT: ')
# -

TwoD_FFT_xr(grid_3D_sg)





data_varis_list = list(grid_3D_gap.variables.mapping.keys())
data_varis_list



# +
kdeplot = grid_3D_sg_pks_pad.LIX_fb_peaks_pad.to_dataframe().plot.kde(ind = 200)
# how to get a line curves. 
kde_data = np.vstack((kdeplot.lines[0]._x, kdeplot.lines[0]._y))
#kde_data.shape
#kde_data[0]

#kde 에서 peak 찾기 
kde_peaks, _= sp.signal.find_peaks(kde_data[1])
kde_properties = sp.signal.peak_widths(kde_data[1], kde_peaks, rel_height=0.5)
# FWHM range  (width, width)height, left_ips, right_ips)

for pk_i in range(len(kde_peaks)):
    print (kde_peaks[pk_i], int( kde_properties[2][pk_i]), int(kde_properties[3][pk_i]))
    # peak 의 갯수에 따라서 , peak 위치 FWHM 의 left& right


#kde_peaks
#[peak poistion]
#kde_properties
#[ [prominence] , [left], [right]]
kde_properties


# peak 1 --- N 
# 

# +
kde_curve = pd.DataFrame(np.array(kdeplot.lines[0]._x, kdeplot.lines[0]._y), columns = ['kde_x','kde_y'])

kde_curve
#
#kde_peaks,kde_properties  = sp.signal.find_peaks(kdeplot)
# -

from scipy.stats import gaussian_kde
sp.stats.gaussian_kde(grid_3D_sg_pks_pad.peaks, , bandwidth=0.2)

# +
            
peaks = grid_3D_sg_pks.LIX_fb_peaks[:3,:4]
#print (peaks)
#peaks.shape
peak_lists = peaks.values

peak_lists_count_max = max([ len(peak_list)
                            for peak_list in peak_lists]
                          )
padding_value = np.nan

peak_lists.shape

peak_lists_count_max = max([ len(peak_lists_r_c) 
                            for peak_lists_r in peak_lists 
                            for  peak_lists_r_c in peak_lists_r])

peak_lists_pad = np.array(
    [np.pad(peak_lists_r_c.astype(float), 
            (0,peak_lists_count_max-len(peak_lists_r_c)),
            mode = 'constant', 
            constant_values = padding_value)
     for peak_lists_r in peak_lists 
     for  peak_lists_r_c in peak_lists_r])


peak_lists_pad.reshape(3,4,-1)

# -

#     * flatten data set  -> filtering with histo widths 
#    * thresholding with peaks --> if the XY position has a peak at certain energy -> ROI 

# dimension size 
grid_3D.dims.mapping['X']
grid_3D.dims.mapping['Y']
grid_3D.dims.mapping['bias_mV']
grid_3D.bias_mV.size

# +
grid_3D.LIX_unit_calc[0,0,:].plot()
plt.show()

grid_3D_sg.LIX_unit_calc[0,0,:].plot()

#grid_3D_2deriv_sg.I_fb[0,0,:].plot()
#plt.show()
#grid_3D_1deriv_sg.LIX_unit_calc[0,0,:].plot()
plt.show()

#grid_3D_sg.differentiate('bias_mV').I_fb[0,0,:].plot()

# +
# position dependent Curves 
# comparison beteween UnitCal

grid_3D_hv = hv.Dataset(grid_3D.LIX_unit_calc).to(hv.Curve,
                                                  kdims =  ["bias_mV"],
                                                  dynamic = True ).relabel('Measured LIX data ')
grid_3D_sg_hv = hv.Dataset(grid_3D_sg.LIX_unit_calc).to(hv.Curve,
                                                  kdims =  ["bias_mV"],
                                                  dynamic = True ).relabel('SG-Smoothing\n of Measured LIX data ')
grid_3D_1deriv_sg_hv = hv.Dataset(grid_3D_1deriv_sg.I_fb).to(hv.Curve,
                                                  kdims =  ["bias_mV"],
                                                  dynamic = True ).relabel('dI/dV after SG-smooth')
grid_3D_hv+grid_3D_sg_hv+grid_3D_1deriv_sg_hv
# -

grid_3D_peaks.I_fb_peaks.values.ravel().

# !pip install tensorflow

# +
# to call the peak list np 
paddingtest= grid_3D_peaks.I_fb_peaks.values


import tensorflow as tf
from tensorflow.keras.preprocessing import sequence

sequence.pad_sequences(paddingtest)




# +

# position dependent Peak Detection 

from scipy.signal import find_peaks, peak_prominences, peak_widths


data_with_peaks = grid_3D_sg.LIX_unit_calc[40,10,:]
peaks,_ = find_peaks (data_with_peaks) 
results_half = peak_widths(data_with_peaks, peaks, rel_height=0.5)
results_full = peak_widths(data_with_peaks, peaks, rel_height=1)

plt.plot(peaks, data_with_peaks[peaks],"xr")
plt.plot(data_with_peaks)
plt.hlines(*results_half[1:], color="C2")
plt.hlines(*results_full[1:], color="C3")

# -

grid_3D_hv = hv.Dataset(grid_3D.LIX_unit_calc).to(hv.Curve,
                                                  kdims =  ["bias_mV"],
                                                  dynamic = True ).relabel('Measured LIX data ')
grid_3D_sg_hv = hv.Dataset(grid_3D_sg.LIX_unit_calc).to(hv.Curve,
                                                  kdims =  ["bias_mV"],
                                                  dynamic = True ).relabel('SG-Smoothing\n of Measured LIX data ')
grid_3D_1deriv_sg_hv = hv.Dataset(grid_3D_1deriv_sg.I_fb).to(hv.Curve,
                                                  kdims =  ["bias_mV"],
                                                  dynamic = True ).relabel('dI/dV after SG-smooth')
grid_3D_hv+grid_3D_sg_hv+grid_3D_1deriv_sg_hv

# +
grid_3D_sg.LIX_unit_calc[0,0,:].plot()
plt.show()
grid_3D_sg.LIX_fb[0,0,:].plot()
v

plt.show()
# -

from holoviews.plotting.links import DataLink
Curve1 = hv.Scatter(grid_3D.isel(X = 1, Y =1).dIdV,  label = 'dIdV')
Curve2 = hv.Scatter(grid_3D_gap.isel(X = 1, Y =1).LDOS_fb,  label = 'grid_3D_gap.LDOS_fb').opts(axiswise=True, ylabel = 'dI/dV [A/V]')
dlink = DataLink(Curve1, Curve2)
# scatter 의 경우만 link 가 표시가능함 
(Curve1+Curve2).opts(opts.Scatter(tools=['box_select', 'lasso_select']))

# +
test_df = grid_3D_sg.to_dataframe()

z_profile  = (PlaneFitY_xr(grid_topo) - PlaneFitY_xr(grid_topo).min()).to_dataframe()
layers = [0,1,2,3,4]
z_profile['z_cut']= pd.cut(z_profile['topography'], labels = layers, bins=5)
z_profile.to_xarray()
